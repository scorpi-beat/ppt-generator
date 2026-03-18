"""
Detailed analysis of template_test1.pptx
"""
import json
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

EMU_PER_PT = 12700

def emu_to_pt(emu):
    if emu is None:
        return None
    return round(emu / EMU_PER_PT, 2)

def shape_type_name(st):
    try:
        return str(MSO_SHAPE_TYPE(st))
    except Exception:
        return str(st)

def placeholder_type_name(ph_idx):
    try:
        return str(PP_PLACEHOLDER(ph_idx))
    except Exception:
        return str(ph_idx)

def get_text_safe(shape):
    try:
        return shape.text.strip()
    except Exception:
        return ""

def looks_like_placeholder_text(text):
    """Detect bracket patterns, Korean labels, or generic English placeholders."""
    import re
    if not text:
        return False
    patterns = [
        r"\[.*?\]",            # [bracketed]
        r"\{.*?\}",            # {curly}
        r"제목|부제목|본문|내용|소제목|설명|텍스트|머리말|꼬리말",  # Korean labels
        r"Title|Subtitle|Body|Content|Header|Footer|Click to edit",  # English
        r"^\d+$",              # lone number
        r"날짜|Date|페이지|Page",
    ]
    for p in patterns:
        if re.search(p, text, re.IGNORECASE):
            return True
    return False

def analyze_shape(shape, idx):
    info = {
        "index": idx,
        "name": shape.name,
        "shape_type": shape_type_name(shape.shape_type),
        "left_emu": shape.left,
        "top_emu": shape.top,
        "width_emu": shape.width,
        "height_emu": shape.height,
        "left_pt": emu_to_pt(shape.left),
        "top_pt": emu_to_pt(shape.top),
        "width_pt": emu_to_pt(shape.width),
        "height_pt": emu_to_pt(shape.height),
        "is_placeholder": shape.is_placeholder,
        "placeholder_type": None,
        "placeholder_idx": None,
        "text": "",
        "text_is_placeholder_like": False,
        "has_table": False,
        "table_info": None,
        "has_chart": False,
        "chart_type": None,
        "fill_type": None,
    }

    # Placeholder details
    if shape.is_placeholder:
        ph = shape.placeholder_format
        info["placeholder_idx"] = ph.idx
        info["placeholder_type"] = placeholder_type_name(ph.type)

    # Text
    text = get_text_safe(shape)
    info["text"] = text[:400]  # cap long text
    info["text_is_placeholder_like"] = looks_like_placeholder_text(text)

    # Table
    if shape.has_table:
        info["has_table"] = True
        tbl = shape.table
        headers = []
        try:
            for cell in tbl.rows[0].cells:
                headers.append(cell.text.strip())
        except Exception:
            pass
        info["table_info"] = {
            "rows": len(tbl.rows),
            "cols": len(tbl.columns),
            "headers": headers,
        }

    # Chart
    if shape.has_chart:
        info["has_chart"] = True
        try:
            info["chart_type"] = str(shape.chart.chart_type)
        except Exception:
            info["chart_type"] = "unknown"

    # Fill type (background color hint)
    try:
        fill = shape.fill
        info["fill_type"] = str(fill.type)
    except Exception:
        pass

    return info

def classify_slide(shapes_info, slide_w_pt, slide_h_pt):
    """Heuristically map to pipeline layout type."""
    texts = [s["text"].lower() for s in shapes_info if s["text"]]
    has_table = any(s["has_table"] for s in shapes_info)
    has_chart = any(s["has_chart"] for s in shapes_info)
    n_shapes = len(shapes_info)
    ph_types = [s["placeholder_type"] for s in shapes_info if s["placeholder_type"]]
    ph_texts = [s["text"] for s in shapes_info if s["text_is_placeholder_like"]]

    # Count text boxes in columns (rough heuristic)
    text_shapes = [s for s in shapes_info if s["text"] or s["is_placeholder"]]
    lefts = sorted(set(round(s["left_pt"] / 10) * 10 for s in text_shapes if s["left_pt"] is not None))

    reasons = []

    if has_table:
        reasons.append("table_slide (has table)")
    if has_chart:
        reasons.append("content_chart (has chart)")

    # Title slide: center-heavy, large title, few shapes
    if n_shapes <= 4 and any("title" in (s["placeholder_type"] or "").lower() for s in shapes_info):
        all_center = all(
            s["left_pt"] is not None and s["left_pt"] > slide_w_pt * 0.1
            for s in text_shapes
        )
        reasons.append("title_slide (few shapes, has title placeholder)")

    # Section divider: very few shapes, large text
    if n_shapes <= 3:
        reasons.append("section_divider (very few shapes)")

    # Two-column: shapes split roughly left/right halves
    left_half = [s for s in text_shapes if s["left_pt"] is not None and s["left_pt"] < slide_w_pt * 0.5]
    right_half = [s for s in text_shapes if s["left_pt"] is not None and s["left_pt"] >= slide_w_pt * 0.5]
    if len(left_half) >= 1 and len(right_half) >= 1 and not has_table and not has_chart:
        reasons.append("two_column_compare (shapes in both halves)")

    # Three-column: multiple shapes spread across thirds
    thirds = [0, slide_w_pt / 3, slide_w_pt * 2 / 3]
    col_counts = [0, 0, 0]
    for s in text_shapes:
        if s["left_pt"] is not None:
            for i in range(2, -1, -1):
                if s["left_pt"] >= thirds[i]:
                    col_counts[i] += 1
                    break
    if all(c >= 1 for c in col_counts):
        reasons.append("three_column_summary (shapes across 3 thirds)")

    # KPI: 4+ shapes arranged in row
    if n_shapes >= 5:
        top_vals = [s["top_pt"] for s in text_shapes if s["top_pt"] is not None]
        if top_vals:
            avg_top = sum(top_vals) / len(top_vals)
            same_row = sum(1 for v in top_vals if abs(v - avg_top) < 50)
            if same_row >= 4:
                reasons.append("kpi_metrics (4+ shapes in same row)")

    # Timeline: many shapes in a horizontal row with roughly equal spacing
    if n_shapes >= 5:
        reasons.append("roadmap_timeline (many shapes, possible timeline)")

    # Default fallback
    if not reasons:
        reasons.append("content_text (fallback)")

    return reasons

def analyze_pptx(path):
    prs = Presentation(path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_w_pt = emu_to_pt(slide_w)
    slide_h_pt = emu_to_pt(slide_h)

    print("=" * 70)
    print(f"FILE: {path}")
    print(f"SLIDE SIZE: {slide_w} x {slide_h} EMU  |  {slide_w_pt} x {slide_h_pt} pt")
    print(f"NUMBER OF SLIDES: {len(prs.slides)}")
    print("=" * 70)

    all_slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'─'*70}")
        print(f"SLIDE {slide_num}")
        print(f"  Layout name: {slide.slide_layout.name if slide.slide_layout else 'N/A'}")
        print(f"  Number of shapes: {len(slide.shapes)}")

        shapes_info = []
        for idx, shape in enumerate(slide.shapes):
            info = analyze_shape(shape, idx)
            shapes_info.append(info)

        # Print shapes
        for s in shapes_info:
            print(f"\n  Shape [{s['index']}] '{s['name']}'")
            print(f"    Type: {s['shape_type']}")
            print(f"    Position: L={s['left_emu']} T={s['top_emu']} W={s['width_emu']} H={s['height_emu']} EMU")
            print(f"             L={s['left_pt']}pt T={s['top_pt']}pt W={s['width_pt']}pt H={s['height_pt']}pt")
            if s["is_placeholder"]:
                print(f"    Placeholder: type={s['placeholder_type']}  idx={s['placeholder_idx']}")
            if s["text"]:
                preview = s["text"][:200].replace('\n', '[NL]')
                print(f"    Text: \"{preview}\"")
                if s["text_is_placeholder_like"]:
                    print(f"    *** PLACEHOLDER-LIKE TEXT DETECTED ***")
            if s["has_table"]:
                ti = s["table_info"]
                print(f"    TABLE: {ti['rows']} rows x {ti['cols']} cols")
                print(f"    Headers: {ti['headers']}")
            if s["has_chart"]:
                print(f"    CHART: type={s['chart_type']}")
            if s["fill_type"] and s["fill_type"] != "None":
                print(f"    Fill: {s['fill_type']}")

        # Classification
        mapping = classify_slide(shapes_info, slide_w_pt, slide_h_pt)
        print(f"\n  >> LAYOUT MAPPING CANDIDATES: {mapping}")

        all_slides.append({
            "slide_num": slide_num,
            "layout_name": slide.slide_layout.name if slide.slide_layout else None,
            "shape_count": len(shapes_info),
            "shapes": shapes_info,
            "layout_candidates": mapping,
        })

    # Summary
    print("\n" + "=" * 70)
    print("LAYOUT MAPPING SUMMARY")
    print("=" * 70)
    for s in all_slides:
        print(f"  Slide {s['slide_num']:2d} | layout='{s['layout_name']}' | shapes={s['shape_count']:2d} | candidates={s['layout_candidates']}")

    return {
        "file": path,
        "slide_width_emu": slide_w,
        "slide_height_emu": slide_h,
        "slide_width_pt": slide_w_pt,
        "slide_height_pt": slide_h_pt,
        "slide_count": len(prs.slides),
        "slides": all_slides,
    }

if __name__ == "__main__":
    import sys, io
    # Redirect stdout to a UTF-8 file so we avoid Windows console encoding issues
    log_path = r"C:\Side project\ppt generator\outputs\template_test1_analysis.txt"
    with open(log_path, "w", encoding="utf-8") as log_file:
        original_stdout = sys.stdout
        sys.stdout = log_file
        result = analyze_pptx(r"C:\Side project\ppt generator\outputs\template_test1.pptx")
        sys.stdout = original_stdout

    # Save JSON summary
    out_path = r"C:\Side project\ppt generator\outputs\template_test1_analysis.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2, default=str)

    sys.stdout = open(sys.stdout.fileno(), mode='w', encoding='utf-8', errors='replace', closefd=False)
    print(f"Analysis complete.")
    print(f"Text log: {log_path}")
    print(f"JSON: {out_path}")
