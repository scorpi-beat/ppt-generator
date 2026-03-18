
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import defaultdict

EMU_PER_CM = 360000

def emu_to_cm(emu):
    return round(emu / EMU_PER_CM, 2)

def get_rgb(color_obj):
    try:
        if color_obj and color_obj.type is not None:
            rgb = color_obj.rgb
            return str(rgb)
    except:
        pass
    return None

def analyze_shape(shape):
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left_cm": emu_to_cm(shape.left) if shape.left is not None else None,
        "top_cm": emu_to_cm(shape.top) if shape.top is not None else None,
        "width_cm": emu_to_cm(shape.width) if shape.width is not None else None,
        "height_cm": emu_to_cm(shape.height) if shape.height is not None else None,
    }

    # Fill color — TABLE/GROUP do not support fill
    info["fill_color"] = None
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            info["fill_color"] = get_rgb(fc)
    except Exception:
        pass

    # Line color
    info["line_color"] = None
    info["line_width_pt"] = None
    try:
        line = shape.line
        info["line_color"] = get_rgb(line.color)
        info["line_width_pt"] = round(line.width / 12700, 1) if line.width else None
    except Exception:
        pass

    # Text
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        info["has_text"] = True
        paras = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run_info = {
                    "text_preview": run.text[:60],
                    "font_name": run.font.name,
                    "font_size_pt": round(run.font.size / 12700, 1) if run.font.size else None,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": get_rgb(run.font.color),
                }
                paras.append(run_info)
        info["runs"] = paras
    else:
        info["has_text"] = False
        info["runs"] = []

    return info

def analyze_slide(slide, idx):
    result = {
        "slide_index": idx + 1,
        "layout_name": slide.slide_layout.name if slide.slide_layout else "None",
        "shapes": [],
        "placeholders": [],
    }

    for ph in slide.placeholders:
        ph_info = {
            "idx": ph.placeholder_format.idx,
            "type": str(ph.placeholder_format.type),
            "name": ph.name,
            "left_cm": emu_to_cm(ph.left) if ph.left is not None else None,
            "top_cm": emu_to_cm(ph.top) if ph.top is not None else None,
            "width_cm": emu_to_cm(ph.width) if ph.width is not None else None,
            "height_cm": emu_to_cm(ph.height) if ph.height is not None else None,
            "text_preview": ph.text[:80] if ph.has_text_frame else "",
        }
        result["placeholders"].append(ph_info)

    for shape in slide.shapes:
        result["shapes"].append(analyze_shape(shape))

    return result

def analyze_master(master, idx):
    result = {
        "master_index": idx,
        "layout_count": len(master.slide_layouts),
        "layout_names": [l.name for l in master.slide_layouts],
        "shape_count": len(master.shapes),
        "shapes": [],
    }
    for shape in master.shapes:
        result["shapes"].append(analyze_shape(shape))
    return result

def collect_colors(prs):
    colors = defaultdict(int)
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fc = get_rgb(shape.fill.fore_color)
                if fc:
                    colors[fc] += 1
            except:
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        c = get_rgb(run.font.color)
                        if c:
                            colors[c] += 1
    return dict(sorted(colors.items(), key=lambda x: -x[1]))

# ---- Load ----
path_b = r"C:\Side project\ppt generator\outputs\report_offshore_wind.pptx"
path_a = r"C:\Side project\ppt generator\outputs\report_offshore_wind_A.pptx"

prs_b = Presentation(path_b)
prs_a = Presentation(path_a)

# ---- Basic Info ----
print("=" * 70)
print("1. 기본 정보")
print("=" * 70)
print(f"{'':30s} {'Plan B (AI)':>18s} {'Plan A (Autofill)':>18s}")
print(f"{'슬라이드 수':30s} {len(prs_b.slides):>18d} {len(prs_a.slides):>18d}")
print(f"{'슬라이드 폭 (cm)':30s} {emu_to_cm(prs_b.slide_width):>18.2f} {emu_to_cm(prs_a.slide_width):>18.2f}")
print(f"{'슬라이드 높이 (cm)':30s} {emu_to_cm(prs_b.slide_height):>18.2f} {emu_to_cm(prs_a.slide_height):>18.2f}")
print(f"{'슬라이드 마스터 수':30s} {len(prs_b.slide_masters):>18d} {len(prs_a.slide_masters):>18d}")

# ---- Master Analysis ----
print("\n" + "=" * 70)
print("2. 슬라이드 마스터 / 레이아웃")
print("=" * 70)
for label, prs in [("Plan B (AI)", prs_b), ("Plan A (Autofill)", prs_a)]:
    print(f"\n  [{label}]")
    for i, master in enumerate(prs.slide_masters):
        m = analyze_master(master, i)
        print(f"    마스터[{i}]: 레이아웃 {m['layout_count']}개, 마스터 도형 {m['shape_count']}개")
        print(f"    레이아웃 이름: {m['layout_names']}")
        shape_types = [s['shape_type'] for s in m['shapes']]
        print(f"    마스터 도형 타입: {shape_types}")

# ---- Slide-by-Slide Analysis (1~5 + last) ----
print("\n" + "=" * 70)
print("3. 슬라이드별 분석 (1~5슬라이드 + 마지막)")
print("=" * 70)

def target_slides(prs):
    n = len(prs.slides)
    indices = list(range(min(5, n)))
    if n > 5:
        indices.append(n - 1)
    return list(dict.fromkeys(indices))

for label, prs in [("Plan B (AI)", prs_b), ("Plan A (Autofill)", prs_a)]:
    print(f"\n{'='*34} {label} {'='*34}")
    for idx in target_slides(prs):
        slide = prs.slides[idx]
        result = analyze_slide(slide, idx)
        total_shapes = len(result["shapes"])
        text_shapes = sum(1 for s in result["shapes"] if s["has_text"])
        non_text_shapes = total_shapes - text_shapes
        ph_count = len(result["placeholders"])

        print(f"\n  [슬라이드 {idx+1}] 레이아웃: {result['layout_name']}")
        print(f"    전체 도형: {total_shapes}개 | 텍스트 포함: {text_shapes}개 | 비텍스트: {non_text_shapes}개 | 플레이스홀더: {ph_count}개")

        # Placeholder info
        for ph in result["placeholders"]:
            print(f"    PH[{ph['idx']}] {ph['type']:30s} pos=({ph['left_cm']},{ph['top_cm']}) size=({ph['width_cm']}x{ph['height_cm']}) text='{ph['text_preview'][:50]}'")

        # Shape detail
        for s in result["shapes"]:
            shape_desc = f"  {s['name'][:30]:30s} type={s['shape_type']:20s} pos=({s['left_cm']},{s['top_cm']}) size=({s['width_cm']}x{s['height_cm']})"
            if s["fill_color"]:
                shape_desc += f" fill=#{s['fill_color']}"
            print("    " + shape_desc)
            if s["has_text"] and s["runs"]:
                for r in s["runs"][:3]:
                    run_desc = f"      → '{r['text_preview'][:50]}' font={r['font_name']} size={r['font_size_pt']}pt bold={r['bold']} color=#{r['color']}"
                    print(run_desc)

# ---- Color Palette ----
print("\n" + "=" * 70)
print("4. 색상 팔레트 (상위 15개)")
print("=" * 70)
for label, prs in [("Plan B (AI)", prs_b), ("Plan A (Autofill)", prs_a)]:
    colors = collect_colors(prs)
    print(f"\n  [{label}]")
    for i, (color, count) in enumerate(list(colors.items())[:15]):
        print(f"    #{color}: {count}회")

# ---- Shape Type Summary ----
print("\n" + "=" * 70)
print("5. 전체 도형 타입 요약")
print("=" * 70)
for label, prs in [("Plan B (AI)", prs_b), ("Plan A (Autofill)", prs_a)]:
    type_counts = defaultdict(int)
    has_image = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            type_counts[str(shape.shape_type)] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image += 1
    print(f"\n  [{label}]")
    for t, c in sorted(type_counts.items(), key=lambda x: -x[1]):
        print(f"    {t}: {c}개")
    print(f"    이미지 도형: {has_image}개")

print("\n분석 완료.")
