"""
extract_spec.py — PPTX 슬라이드 외부(pasteboard) 스펙 자동 추출

슬라이드 마스터의 캔버스 바깥 영역에 작성된 규격(폰트 크기, 컬러 팔레트,
표 스타일)을 읽어 outputs/spec_{stem}.json 으로 저장한다.

사용:
  py src/core/extract_spec.py outputs/template_test1.pptx

이후 assembler.py / zone_layout.py 는 이 JSON을 자동 로드해 적용한다.
파이프라인 실행 전 이 스크립트를 먼저 실행하거나,
assembler.py 가 JSON 없을 때 자동 실행한다.
"""

import argparse
import json
import os
import re
import sys

# ─────────────────────────────────────────────────────────────────────────────
# 한국어 폰트 크기 레이블 → spec 키 매핑
# ─────────────────────────────────────────────────────────────────────────────
FONT_LABEL_MAP = {
    "제목":     "title_pt",
    "중제목":   "sub_title_pt",
    "소제목":   "section_title_pt",
    "본문":     "body_pt",
    "각주":     "footnote_pt",
    "출처":     "footnote_pt",
    "표 구분":  "table_header_pt",
    "표구분":   "table_header_pt",
    "표 내용":  "table_body_pt",
    "표내용":   "table_body_pt",
    "줄간격":   "line_spacing",
}

# 컬러 스와치 x 좌표 → 레이블 텍스트 근접 매칭 허용 오차 (pt)
LABEL_MATCH_TOLERANCE_PT = 50


# ─────────────────────────────────────────────────────────────────────────────
# 폰트 크기 파싱
# ─────────────────────────────────────────────────────────────────────────────

def _parse_font_sizes(text: str) -> dict:
    """
    "제목 27p / 중제목 18p / 소제목 16p / 본문 14p / 각주/출처 9p
     표 구분 14p / 표 내용 12p / 줄간격 1.5"
    형식의 문자열에서 폰트 크기를 파싱한다.
    """
    result = {}
    # 슬래시, 줄바꿈, 탭으로 분리
    tokens = re.split(r"[/\n\t]", text)
    for token in tokens:
        token = token.strip()
        # "레이블 숫자p" 또는 "레이블 숫자" 패턴
        m = re.search(r"([\w\s/]+?)\s+(\d+(?:\.\d+)?)p?$", token, re.UNICODE)
        if not m:
            continue
        label_raw = m.group(1).strip()
        value = float(m.group(2))
        # 레이블 매핑
        for label_key, spec_key in FONT_LABEL_MAP.items():
            if label_key in label_raw:
                if spec_key not in result:
                    result[spec_key] = value
                break
    return result


# ─────────────────────────────────────────────────────────────────────────────
# 표 스타일 파싱
# ─────────────────────────────────────────────────────────────────────────────

def _parse_table_style(text: str, palette_colors: list) -> dict:
    """
    "3. 위쪽/아래쪽 테두리, 1pt, 펜색 황갈색(팔레트 4번째)
     4. 색 채운 부분 안쪽 테두리, 0.5pt, 펜색 흰색
     5. 색 안채운 부분 안쪽 테두리, 0.5pt, 펜색 황갈색" 파싱
    palette_colors: 스와치에서 추출한 색상 리스트 (순서 유지, 1-indexed)
    """
    style = {}

    def _extract_width(line: str) -> float | None:
        m = re.search(r"(\d+(?:\.\d+)?)pt", line)
        return float(m.group(1)) if m else None

    def _resolve_color(line: str, palette: list) -> str | None:
        # "팔레트 N번째" 패턴
        m = re.search(r"팔레트\s*(\d+)번째", line)
        if m:
            idx = int(m.group(1)) - 1   # 1-indexed → 0-indexed
            if 0 <= idx < len(palette):
                return palette[idx]
        if "흰색" in line or "white" in line.lower():
            return "#FFFFFF"
        if "검정" in line or "black" in line.lower():
            return "#000000"
        return None

    lines = text.split("\n")
    for line in lines:
        line = line.strip()
        if re.search(r"3[.\)]\s*위", line) or "위쪽" in line and "아래쪽" in line:
            style["outer_border_width_pt"]  = _extract_width(line)
            style["outer_border_color"]     = _resolve_color(line, palette_colors)
        elif re.search(r"4[.\)]\s*색", line) or ("채운" in line and "안쪽" in line and "색 안" not in line):
            style["inner_filled_width_pt"]  = _extract_width(line)
            style["inner_filled_color"]     = _resolve_color(line, palette_colors)
        elif re.search(r"5[.\)]\s*색", line) or ("안채운" in line and "안쪽" in line) or ("색 안" in line and "안쪽" in line):
            style["inner_unfilled_width_pt"] = _extract_width(line)
            style["inner_unfilled_color"]    = _resolve_color(line, palette_colors)

    return style


# ─────────────────────────────────────────────────────────────────────────────
# 컬러 스와치 추출
# ─────────────────────────────────────────────────────────────────────────────

def _hex_from_shape(shape) -> str | None:
    """도형의 solid fill 색상을 HEX 문자열로 반환."""
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.util import Pt
        fore = fill.fore_color
        # RGB 직접 접근
        rgb = fore.rgb
        return f"#{rgb}"
    except Exception:
        return None


def _extract_palette(master, slide_w_pt: float, slide_h_pt: float) -> tuple[list, dict]:
    """
    마스터 pasteboard 에서 컬러 스와치와 레이블을 추출한다.
    반환: (colors_in_order, {role_label: hex_color})
    """
    from pptx.util import Pt
    EMU = 12700

    swatches = []   # (left_pt, hex_color)
    labels   = []   # (left_pt, label_text)

    for shape in master.shapes:
        left_pt = shape.left  / EMU
        top_pt  = shape.top   / EMU

        # pasteboard 판별 (슬라이드 영역 밖)
        is_offcanvas = (top_pt < 0 or left_pt < 0
                        or top_pt > slide_h_pt or left_pt > slide_w_pt)
        if not is_offcanvas:
            continue

        # 색상 스와치: fill이 있는 AUTO_SHAPE
        if shape.shape_type in (1, 5):   # AUTO_SHAPE
            hex_color = _hex_from_shape(shape)
            if hex_color and hex_color not in ("#FFFFFF", "#000000"):
                swatches.append((left_pt, hex_color))

        # 레이블 텍스트박스
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt and len(txt) < 30 and "{{" not in txt:
                labels.append((left_pt, txt))

    # x 좌표 순으로 정렬
    swatches.sort(key=lambda x: x[0])
    colors_in_order = [c for _, c in swatches]

    # 레이블-스와치 매핑: x 좌표 근접 매칭
    role_map = {}
    for label_x, label_text in labels:
        best_dist = LABEL_MATCH_TOLERANCE_PT
        best_color = None
        for swatch_x, hex_color in swatches:
            dist = abs(swatch_x - label_x)
            if dist < best_dist:
                best_dist = dist
                best_color = hex_color
        if best_color:
            # 레이블 정규화
            key = (label_text
                   .replace(",", "")
                   .replace("/", "_")
                   .strip()
                   .lower()
                   .replace(" ", "_"))
            role_map[key] = best_color

    return colors_in_order, role_map


# ─────────────────────────────────────────────────────────────────────────────
# 메인 추출 함수
# ─────────────────────────────────────────────────────────────────────────────

def extract_spec(pptx_path: str, output_path: str | None = None) -> dict:
    """
    PPTX pasteboard 에서 스펙을 추출해 dict 로 반환하고 JSON 파일로 저장.
    output_path 미지정 시 outputs/spec_{stem}.json 에 저장.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("[오류] python-pptx 미설치. pip install python-pptx 실행 후 재시도.")
        sys.exit(1)

    EMU = 12700

    prs       = Presentation(pptx_path)
    slide_w   = prs.slide_width  / EMU   # pt
    slide_h   = prs.slide_height / EMU   # pt
    master    = prs.slide_master

    spec = {
        "source_file": os.path.basename(pptx_path),
        "slide_size": {
            "width_emu":  int(prs.slide_width),
            "height_emu": int(prs.slide_height),
            "width_pt":   round(slide_w, 2),
            "height_pt":  round(slide_h, 2),
        },
        "font_sizes": {
            "title_pt":         27,
            "sub_title_pt":     18,
            "section_title_pt": 16,
            "body_pt":          14,
            "footnote_pt":       9,
            "table_header_pt":  14,
            "table_body_pt":    12,
            "line_spacing":      1.5,
        },
        "color_palette": {},
        "color_roles": {},
        "table_style": {
            "outer_border_width_pt":  1.0,
            "outer_border_color":    "#BCB8AF",
            "inner_filled_width_pt":  0.5,
            "inner_filled_color":    "#FFFFFF",
            "inner_unfilled_width_pt": 0.5,
            "inner_unfilled_color":  "#BCB8AF",
            "header_fill":           "#2D3734",
            "header_font_color":     "#FFFFFF",
            "row_fill_even":         "#E4E0D4",
            "row_fill_odd":          "#FFFFFF",
        },
    }

    # ── 1. pasteboard 텍스트에서 폰트 크기 + 표 스타일 추출 ──────────────
    font_text_found = ""
    table_style_text = ""

    for shape in master.shapes:
        left_pt = shape.left / EMU
        top_pt  = shape.top  / EMU
        is_off  = (top_pt < 0 or left_pt < 0
                   or top_pt > slide_h or left_pt > slide_w)
        if not is_off or not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if not txt:
            continue

        # 폰트 크기 가이드 판별 (p 단위 숫자 포함)
        if re.search(r"\d+p\b", txt) and ("제목" in txt or "본문" in txt):
            font_text_found = txt

        # 표 스타일 가이드 판별 (pt 단위 + 테두리 언급)
        if "pt" in txt and ("테두리" in txt or "채우기" in txt):
            table_style_text = txt

    if font_text_found:
        parsed_fonts = _parse_font_sizes(font_text_found)
        spec["font_sizes"].update(parsed_fonts)
        print(f"  [폰트] 추출 완료: {parsed_fonts}")

    # ── 2. 컬러 스와치 추출 ───────────────────────────────────────────────
    colors_ordered, role_map = _extract_palette(master, slide_w, slide_h)
    spec["color_palette"] = {f"color_{i+1}": c for i, c in enumerate(colors_ordered)}
    spec["color_roles"]   = role_map

    if colors_ordered:
        print(f"  [팔레트] {len(colors_ordered)}개 색상 추출: {colors_ordered}")

    # ── 3. 표 스타일 파싱 ─────────────────────────────────────────────────
    if table_style_text:
        parsed_style = _parse_table_style(table_style_text, colors_ordered)
        spec["table_style"].update({k: v for k, v in parsed_style.items() if v is not None})
        print(f"  [표 스타일] 추출 완료: {parsed_style}")

    # ── 4. JSON 저장 ──────────────────────────────────────────────────────
    if output_path is None:
        stem = os.path.splitext(os.path.basename(pptx_path))[0]
        base_dir = os.path.normpath(
            os.path.join(os.path.dirname(__file__), "..", "..")
        )
        output_path = os.path.join(base_dir, "outputs", f"spec_{stem}.json")

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(spec, f, ensure_ascii=False, indent=2)

    print(f"  스펙 저장 완료: {output_path}")
    return spec


def load_spec(pptx_path: str) -> dict:
    """
    spec JSON 을 로드한다. 없으면 extract_spec() 을 자동 실행한다.
    assembler.py / zone_layout.py 에서 호출.
    """
    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    base_dir = os.path.normpath(
        os.path.join(os.path.dirname(__file__), "..", "..")
    )
    spec_path = os.path.join(base_dir, "outputs", f"spec_{stem}.json")

    if not os.path.exists(spec_path):
        print(f"  [스펙] {spec_path} 없음 → 자동 추출 실행")
        return extract_spec(pptx_path, spec_path)

    with open(spec_path, encoding="utf-8") as f:
        return json.load(f)


# ─────────────────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX pasteboard 스펙 추출")
    parser.add_argument("pptx", help="대상 PPTX 파일 경로")
    parser.add_argument("--out", help="출력 JSON 경로 (기본: outputs/spec_{stem}.json)")
    args = parser.parse_args()

    spec = extract_spec(args.pptx, args.out)
    print(json.dumps(spec, ensure_ascii=False, indent=2))
