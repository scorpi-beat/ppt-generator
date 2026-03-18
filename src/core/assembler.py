"""
assembler.py — Plan A+B 하이브리드 조립기

draft.json + component_template.pptx → 완성 PPTX

사용법:
  python src/core/assembler.py \
    --draft outputs/draft_report_offshore_wind.json \
    --template outputs/component_template_report.pptx \
    --output outputs/report_offshore_wind_hybrid.pptx
"""

import argparse
import copy
import json
import os
import re
import shutil
import sys
import tempfile
import zipfile
from lxml import etree

# 로컬 모듈 경로 등록
sys.path.insert(0, os.path.join(os.path.dirname(__file__)))
from zone_layout import (
    ZONE_CONFIGS, get_zone_rects, load_calibration,
    should_split_slide, SOURCE_Y, SOURCE_CY,
)
from extract_spec import load_spec

# ---------------------------------------------------------------------------
# 텍스트 교체 매핑
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# 토큰 정규화: template_test1.pptx의 비표준 토큰을 표준 draft 키로 정규화.
# replace_placeholders() 가 PLACEHOLDER_MAP 적용 전에 먼저 이 변환을 수행한다.
# ─ 공백 포함 토큰, 대소문자 불일치, 오타 등을 일괄 처리
# ---------------------------------------------------------------------------
TOKEN_NORMALIZE = {
    # 대소문자
    "{{Head Message}}":          "{{head_message}}",
    "{{head message}}":          "{{head_message}}",
    # 공백 → 언더스코어
    "{{section number}}":        "{{section_number}}",
    "{{section name}}":          "{{section_name}}",
    "{{section title}}":         "{{section_title}}",
    "{{chart_02 title}}":        "{{chart_02_title}}",
    "{{name list}}":             "{{name_list}}",
    # 오타
    "{{Data_Refernence}}":       "{{source}}",
    "{{mail_adress}}":           "{{mail_address}}",
}


PLACEHOLDER_MAP = {
    # ── title_slide (슬라이드 1) ──────────────────────────────────────────
    "title_slide": {
        "{{title}}":      lambda s: s.get("title", ""),
        "{{subtitle}}":   lambda s: s.get("subtitle", ""),
        "{{company}}":    lambda s: s.get("company", ""),
        "{{series}}":     lambda s: s.get("series", ""),
        "{{name_list}}":  lambda s: s.get("name_list", ""),
        "{{date}}":       lambda s: s.get("date", ""),
    },
    # ── content_text / zone_base (슬라이드 2) ────────────────────────────
    "content_text": {
        "{{section_number}}": lambda s: s.get("section_number", ""),
        "{{section_name}}":   lambda s: s.get("section_name", s.get("title", "")),
        "{{section_title}}":  lambda s: s.get("section_title", s.get("title", "")),
        "{{head_message}}":   lambda s: s.get("head_message", ""),
        "{{toc}}":            lambda s: s.get("toc", ""),
        "{{toc_1}}":          lambda s: _toc(s, 0),
        "{{toc_2}}":          lambda s: _toc(s, 1),
        "{{toc_03}}":         lambda s: _toc(s, 2),
        "{{source}}":         lambda s: s.get("source", ""),
    },
    "zone_base": {
        "{{section_number}}": lambda s: s.get("section_number", ""),
        "{{section_name}}":   lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":   lambda s: s.get("head_message", ""),
        "{{source}}":         lambda s: s.get("source", ""),
    },
    # ── section_divider (슬라이드 3) ─────────────────────────────────────
    "section_divider": {
        "{{section_number}}": lambda s: s.get("section_number", ""),
        "{{section_title}}":  lambda s: s.get("title", ""),
        "{{section_subtitle}}": lambda s: s.get("subtitle", ""),
        "{{toc_1}}":          lambda s: _toc(s, 0),
        "{{toc_2}}":          lambda s: _toc(s, 1),
        "{{toc_03}}":         lambda s: _toc(s, 2),
    },
    # ── dual_chart / content_chart (슬라이드 4, 5) ───────────────────────
    "dual_chart": {
        "{{section_number}}":  lambda s: s.get("section_number", ""),
        "{{section_name}}":    lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":    lambda s: s.get("head_message", ""),
        "{{chart_01_title}}":  lambda s: _chart(s, 0, "title"),
        "{{chart_02_title}}":  lambda s: _chart(s, 1, "title"),
        "{{source}}":          lambda s: s.get("source", ""),
    },
    "content_chart": {
        "{{section_number}}":  lambda s: s.get("section_number", ""),
        "{{section_name}}":    lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":    lambda s: s.get("head_message", ""),
        "{{chart_label}}":     lambda s: s.get("chart", {}).get("title", ""),
        "{{chart_01_title}}":  lambda s: _chart(s, 0, "title"),
        "{{chart_02_title}}":  lambda s: _chart(s, 1, "title"),
        "{{source}}":          lambda s: s.get("source", ""),
    },
    # ── roadmap_timeline (슬라이드 6) ────────────────────────────────────
    "roadmap_timeline": {
        "{{section_number}}":  lambda s: s.get("section_number", ""),
        "{{section_name}}":    lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":    lambda s: s.get("head_message", ""),
        "{{text_title}}":      lambda s: s.get("text_title", s.get("title", "")),
        "{{header_1}}":        lambda s: _tbl_header(s, 0),
        "{{row1_col1}}":       lambda s: _tbl_cell(s, 0, 0),
        "{{row1_col2}}":       lambda s: _tbl_cell(s, 0, 1),
        "{{icon}}":            lambda s: s.get("icon", ""),
        "{{node_title_1}}":    lambda s: _timeline(s, 0, "title"),
        "{{node_title_2}}":    lambda s: _timeline(s, 1, "title"),
        "{{node_title_3}}":    lambda s: _timeline(s, 2, "title"),
        "{{node_title_4}}":    lambda s: _timeline(s, 3, "title"),
        "{{node_title_5}}":    lambda s: _timeline(s, 4, "title"),
        "{{node_title_6}}":    lambda s: _timeline(s, 5, "title"),
        "{{node_text_1}}":     lambda s: _timeline(s, 0, "description"),
        "{{node_text_2}}":     lambda s: _timeline(s, 1, "description"),
        "{{node_text_3}}":     lambda s: _timeline(s, 2, "description"),
        "{{node_text_4}}":     lambda s: _timeline(s, 3, "description"),
        "{{node_text_5}}":     lambda s: _timeline(s, 4, "description"),
        "{{node_text_6}}":     lambda s: _timeline(s, 5, "description"),
    },
    # ── table_slide (슬라이드 7) ──────────────────────────────────────────
    "table_slide": {
        "{{section_number}}": lambda s: s.get("section_number", ""),
        "{{section_name}}":   lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":   lambda s: s.get("head_message", ""),
        "{{source}}":         lambda s: s.get("source", ""),
    },
    # ── two_column_compare (슬라이드 8) ──────────────────────────────────
    "two_column_compare": {
        "{{section_number}}": lambda s: s.get("section_number", ""),
        "{{section_name}}":   lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":   lambda s: s.get("head_message", ""),
        "{{left_title}}":     lambda s: _col(s, "column_left",  "title"),
        "{{right_title}}":    lambda s: _col(s, "column_right", "title"),
        "{{image_caption}}":  lambda s: s.get("image_caption", ""),
        "{{source}}":         lambda s: s.get("source", ""),
    },
    # ── photo_table (슬라이드 9) ──────────────────────────────────────────
    "photo_table": {
        "{{section_number}}": lambda s: s.get("section_number", ""),
        "{{section_name}}":   lambda s: s.get("section_name", s.get("title", "")),
        "{{head_message}}":   lambda s: s.get("head_message", ""),
        "{{left_title}}":     lambda s: _col(s, "column_left",  "title"),
        "{{source}}":         lambda s: s.get("source", ""),
    },
    # ── image_zones (슬라이드 10) ─────────────────────────────────────────
    "image_zones": {
        "{{section_number}}":     lambda s: s.get("section_number", ""),
        "{{section_name}}":       lambda s: s.get("section_name", s.get("title", "")),
        "{{index_head_message}}": lambda s: _zone_title(s, 0),
        "{{head_message}}":       lambda s: _zone_body_head(s, 0),
    },
    # ── closing_slide (슬라이드 11) ───────────────────────────────────────
    "closing_slide": {
        "{{team_name}}":     lambda s: s.get("team_name", s.get("title", "")),
        "{{name}}":          lambda s: s.get("members", [{}])[0].get("name", ""),
        "{{company}}":       lambda s: s.get("members", [{}])[0].get("company", ""),
        "{{university}}":    lambda s: s.get("members", [{}])[0].get("university", ""),
        "{{major}}":         lambda s: s.get("members", [{}])[0].get("major", ""),
        "{{mail_address}}":  lambda s: s.get("members", [{}])[0].get("email", ""),
        "{{closing_title}}": lambda s: s.get("title", ""),
        "{{closing_message}}": lambda s: s.get("closing_message", ""),
    },
    # ── 기존 호환성 유지 ──────────────────────────────────────────────────
    "three_column_summary": {
        "{{section}}":      lambda s: s.get("title", ""),
        "{{head_message}}": lambda s: s.get("head_message", ""),
        "{{col1_title}}":   lambda s: _col_list(s, 0, "title"),
        "{{col2_title}}":   lambda s: _col_list(s, 1, "title"),
        "{{col3_title}}":   lambda s: _col_list(s, 2, "title"),
    },
    "kpi_metrics": {
        "{{section}}":      lambda s: s.get("title", ""),
        "{{head_message}}": lambda s: s.get("head_message", ""),
        "{{kpi_1_value}}":  lambda s: _kpi(s, 0, "value"),
        "{{kpi_1_label}}":  lambda s: _kpi(s, 0, "label"),
        "{{kpi_1_note}}":   lambda s: _kpi(s, 0, "note"),
        "{{kpi_2_value}}":  lambda s: _kpi(s, 1, "value"),
        "{{kpi_2_label}}":  lambda s: _kpi(s, 1, "label"),
        "{{kpi_2_note}}":   lambda s: _kpi(s, 1, "note"),
        "{{kpi_3_value}}":  lambda s: _kpi(s, 2, "value"),
        "{{kpi_3_label}}":  lambda s: _kpi(s, 2, "label"),
        "{{kpi_3_note}}":   lambda s: _kpi(s, 2, "note"),
        "{{kpi_4_value}}":  lambda s: _kpi(s, 3, "value"),
        "{{kpi_4_label}}":  lambda s: _kpi(s, 3, "label"),
        "{{kpi_4_note}}":   lambda s: _kpi(s, 3, "note"),
    },
    "table_chart_combo": {
        "{{section}}":      lambda s: s.get("title", ""),
        "{{head_message}}": lambda s: s.get("head_message", ""),
        "{{chart_label}}":  lambda s: s.get("chart", {}).get("title", ""),
        "{{header_1}}":     lambda s: _tbl_header(s, 0),
        "{{header_2}}":     lambda s: _tbl_header(s, 1),
        "{{header_3}}":     lambda s: _tbl_header(s, 2),
        "{{row1_col1}}":    lambda s: _tbl_cell(s, 0, 0),
        "{{row1_col2}}":    lambda s: _tbl_cell(s, 0, 1),
        "{{row1_col3}}":    lambda s: _tbl_cell(s, 0, 2),
        "{{row2_col1}}":    lambda s: _tbl_cell(s, 1, 0),
        "{{row2_col2}}":    lambda s: _tbl_cell(s, 1, 1),
        "{{row2_col3}}":    lambda s: _tbl_cell(s, 1, 2),
        "{{row3_col1}}":    lambda s: _tbl_cell(s, 2, 0),
        "{{row3_col2}}":    lambda s: _tbl_cell(s, 2, 1),
        "{{row3_col3}}":    lambda s: _tbl_cell(s, 2, 2),
    },
    "image_gallery": {
        "{{section}}":      lambda s: s.get("title", ""),
        "{{head_message}}": lambda s: s.get("head_message", ""),
        "{{img_1_label}}":  lambda s: _img(s, 0, "label"),
        "{{img_2_label}}":  lambda s: _img(s, 1, "label"),
        "{{img_3_label}}":  lambda s: _img(s, 2, "label"),
        "{{img_4_label}}":  lambda s: _img(s, 3, "label"),
        "{{img_5_label}}":  lambda s: _img(s, 4, "label"),
        "{{img_6_label}}":  lambda s: _img(s, 5, "label"),
    },
}
# zone 레이아웃도 zone_base와 동일한 플레이스홀더 매핑 사용
PLACEHOLDER_MAP["zone"] = PLACEHOLDER_MAP["zone_base"]

# ---------------------------------------------------------------------------
# 데이터 추출 헬퍼
# ---------------------------------------------------------------------------

def _body(slide_data: dict, idx: int) -> str:
    body = slide_data.get("body", [])
    return body[idx] if len(body) > idx else ""


def _key_point(slide_data: dict, idx: int) -> str:
    kp = slide_data.get("key_points", [])
    return kp[idx] if len(kp) > idx else ""


def _col(slide_data: dict, col_key: str, field: str) -> str:
    return slide_data.get(col_key, {}).get(field, "")


def _col_item(slide_data: dict, col_key: str, idx: int) -> str:
    items = slide_data.get(col_key, {}).get("items", [])
    return items[idx] if len(items) > idx else ""


def _col_list(slide_data: dict, col_idx: int, field: str) -> str:
    cols = slide_data.get("columns", [])
    if len(cols) > col_idx:
        return cols[col_idx].get(field, "")
    return ""


def _col_list_item(slide_data: dict, col_idx: int, item_idx: int) -> str:
    cols = slide_data.get("columns", [])
    if len(cols) > col_idx:
        items = cols[col_idx].get("items", [])
        return items[item_idx] if len(items) > item_idx else ""
    return ""


def _tbl_header(slide_data: dict, idx: int) -> str:
    hdrs = slide_data.get("table", {}).get("headers", [])
    return hdrs[idx] if len(hdrs) > idx else ""


def _tbl_cell(slide_data: dict, row_idx: int, col_idx: int) -> str:
    rows = slide_data.get("table", {}).get("rows", [])
    if len(rows) > row_idx:
        row = rows[row_idx]
        return row[col_idx] if len(row) > col_idx else ""
    return ""


def _timeline(slide_data: dict, idx: int, field: str) -> str:
    tl = slide_data.get("timeline", [])
    if len(tl) > idx:
        return tl[idx].get(field, "")
    return ""


def _takeaway(slide_data: dict, idx: int) -> str:
    kts = slide_data.get("key_takeaways", [])
    return kts[idx] if len(kts) > idx else ""


def _kpi(slide_data: dict, idx: int, field: str) -> str:
    kpis = slide_data.get("kpis", [])
    if len(kpis) > idx:
        return str(kpis[idx].get(field, ""))
    return ""


def _img(slide_data: dict, idx: int, field: str) -> str:
    imgs = slide_data.get("images", [])
    if len(imgs) > idx:
        return str(imgs[idx].get(field, ""))
    return ""


def _toc(slide_data: dict, idx: int) -> str:
    """toc 배열에서 idx번째 항목 반환."""
    toc = slide_data.get("toc", [])
    if isinstance(toc, list) and len(toc) > idx:
        return str(toc[idx])
    return ""


def _chart(slide_data: dict, chart_idx: int, field: str) -> str:
    """charts 배열 또는 단일 chart에서 필드 반환."""
    charts = slide_data.get("charts", [])
    if isinstance(charts, list) and len(charts) > chart_idx:
        return str(charts[chart_idx].get(field, ""))
    # 단일 chart 키 (chart_idx==0 인 경우 fallback)
    if chart_idx == 0:
        return str(slide_data.get("chart", {}).get(field, ""))
    return ""


def _zone_title(slide_data: dict, zone_idx: int) -> str:
    """zones[zone_idx].title 반환."""
    zones = slide_data.get("zones", [])
    if len(zones) > zone_idx:
        return str(zones[zone_idx].get("title", ""))
    return ""


def _zone_body_head(slide_data: dict, zone_idx: int) -> str:
    """zones[zone_idx].head_message 또는 첫 번째 body 항목 반환."""
    zones = slide_data.get("zones", [])
    if len(zones) > zone_idx:
        z = zones[zone_idx]
        if z.get("head_message"):
            return str(z["head_message"])
        body = z.get("body", [])
        if isinstance(body, list) and body:
            return str(body[0])
    return ""


def _auto_resolve(key: str, slide_data: dict):
    """
    점 표기법(dot notation)으로 draft 데이터 자동 탐색.
    "chart.title"   → slide_data["chart"]["title"]
    "columns.0.title" → slide_data["columns"][0]["title"]
    리스트인 경우 그대로 반환 (Stage 5 동적 확장 대상).
    """
    parts = key.split(".")
    val = slide_data
    for part in parts:
        if val is None:
            break
        if isinstance(val, dict):
            val = val.get(part)
        elif isinstance(val, list):
            try:
                val = val[int(part)]
            except (ValueError, IndexError):
                val = None
        else:
            val = None
    return val


# ---------------------------------------------------------------------------
# 텍스트 교체
# ---------------------------------------------------------------------------

def _xml_escape(value: str) -> str:
    return (
        value
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


def replace_placeholders(slide_xml_str: str, layout: str, slide_data: dict) -> str:
    """
    1단계: PLACEHOLDER_MAP 명시 항목 교체 (복잡한 경로 표현 포함)
    2단계: 잔여 {{marker}} 자동 탐색 — 문자열/숫자면 즉시 교체, 리스트면 Stage 5로 남김
    """
    # 0단계: 토큰 정규화 (공백/대소문자/오타 → 표준 키)
    for raw, normalized in TOKEN_NORMALIZE.items():
        slide_xml_str = slide_xml_str.replace(raw, normalized)

    # 1단계: 명시 매핑
    mapping = PLACEHOLDER_MAP.get(layout, {})
    for key, value_fn in mapping.items():
        value = value_fn(slide_data)
        slide_xml_str = slide_xml_str.replace(key, _xml_escape(value))

    # 2단계: 자동 탐색 (문자열/숫자 타입만 교체, 리스트는 Stage 5 담당)
    remaining = re.findall(r'\{\{[a-z][a-z0-9_.]*\}\}', slide_xml_str)
    for marker in set(remaining):
        key = marker[2:-2]
        val = _auto_resolve(key, slide_data)
        if isinstance(val, (str, int, float)):
            slide_xml_str = slide_xml_str.replace(marker, _xml_escape(str(val)))
        # list → Stage 5에서 처리, None → 빈 문자열로
        elif val is None:
            slide_xml_str = slide_xml_str.replace(marker, "")

    return slide_xml_str


# ---------------------------------------------------------------------------
# 슬라이드 외부 요소 필터
# ---------------------------------------------------------------------------

def filter_outside_elements(slide_xml_str: str, slide_w: int = 17_610_138, slide_h: int = 9_906_000) -> str:
    """음수 좌표 또는 슬라이드 영역 밖 요소(sp/graphicFrame/grpSp)를 제거."""
    try:
        root = etree.fromstring(slide_xml_str.encode("utf-8"))
    except etree.XMLSyntaxError:
        return slide_xml_str

    ns = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
    spTree_tags = [
        "{http://schemas.openxmlformats.org/presentationml/2006/main}spTree",
    ]

    # a:off 요소를 찾아 x, y 체크
    OFF_TAG = "{http://schemas.openxmlformats.org/drawingml/2006/main}off"
    CANDIDATE_TAGS = {
        "{http://schemas.openxmlformats.org/presentationml/2006/main}sp",
        "{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame",
        "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp",
        "{http://schemas.openxmlformats.org/presentationml/2006/main}pic",
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp",
    }

    def _check_element(elem):
        for off in elem.iter(OFF_TAG):
            try:
                x = int(off.get("x", "0"))
                y = int(off.get("y", "0"))
                if x < 0 or y < 0:
                    return False
            except ValueError:
                pass
        return True

    for spTree in root.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}spTree"):
        to_remove = []
        for child in list(spTree):
            if child.tag in CANDIDATE_TAGS:
                if not _check_element(child):
                    to_remove.append(child)
        for elem in to_remove:
            spTree.remove(elem)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True).decode("utf-8")


# ---------------------------------------------------------------------------
# PPTX ZIP 기반 조립
# ---------------------------------------------------------------------------

def _slide_filename(idx_1based: int) -> str:
    return f"ppt/slides/slide{idx_1based}.xml"


def _slide_rels_filename(idx_1based: int) -> str:
    return f"ppt/slides/_rels/slide{idx_1based}.xml.rels"


def _read_zip_text(zf: zipfile.ZipFile, name: str) -> str:
    with zf.open(name) as f:
        return f.read().decode("utf-8")


def _list_slides(zf: zipfile.ZipFile) -> list:
    """template PPTX 내의 슬라이드 파일 목록을 번호 순으로 반환."""
    names = [n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)]
    names.sort(key=lambda n: int(re.search(r"\d+", os.path.basename(n)).group()))
    return names


def _update_content_types(content_types_xml: str, slide_count: int) -> str:
    """[Content_Types].xml 에 새 슬라이드 Override 항목 추가."""
    root = etree.fromstring(content_types_xml.encode("utf-8"))
    NS = "http://schemas.openxmlformats.org/package/2006/content-types"
    SLIDE_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"

    # 기존 슬라이드 Override 제거 후 재추가
    # notesSlide / notesMaster 항목도 함께 제거 (파일 자체를 삭제했으므로)
    REMOVE_KEYWORDS = ("presentationml.slide+xml", "presentationml.notesSlide", "presentationml.notesMaster")
    to_remove = [
        el for el in root
        if any(kw in (el.get("ContentType") or "") for kw in REMOVE_KEYWORDS)
        or any(kw in (el.get("PartName") or "") for kw in ("notesSlide", "notesMaster"))
    ]
    for el in to_remove:
        root.remove(el)

    for i in range(1, slide_count + 1):
        override = etree.SubElement(root, f"{{{NS}}}Override")
        override.set("PartName", f"/ppt/slides/slide{i}.xml")
        override.set("ContentType", SLIDE_TYPE)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True).decode("utf-8")


def _update_presentation_xml(prs_xml: str, slide_count: int,
                              slide_rids: list[str] | None = None,
                              base_id: int = 256) -> str:
    """ppt/presentation.xml 의 sldIdLst 를 새 슬라이드 목록으로 교체."""
    root = etree.fromstring(prs_xml.encode("utf-8"))
    NS_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    sldIdLst = root.find(f"{{{NS_PML}}}sldIdLst")
    if sldIdLst is None:
        sldIdLst = etree.SubElement(root, f"{{{NS_PML}}}sldIdLst")

    for child in list(sldIdLst):
        sldIdLst.remove(child)

    for i in range(1, slide_count + 1):
        rid = slide_rids[i - 1] if slide_rids else f"rId{i}"
        sldId = etree.SubElement(sldIdLst, f"{{{NS_PML}}}sldId")
        sldId.set("id", str(base_id + i))
        sldId.set(f"{{{NS_R}}}id", rid)

    # notesMasterIdLst 제거
    nml = root.find(f"{{{NS_PML}}}notesMasterIdLst")
    if nml is not None:
        root.remove(nml)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True).decode("utf-8")


def _update_doc_props(app_xml: str, slide_count: int) -> str:
    """docProps/app.xml 의 Slides / Notes 수를 실제 슬라이드 수로 업데이트."""
    NS = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
    try:
        root = etree.fromstring(app_xml.encode("utf-8"))
    except etree.XMLSyntaxError:
        return app_xml
    for tag in ("Slides", "Notes"):
        el = root.find(f"{{{NS}}}{tag}")
        if el is not None:
            el.text = str(slide_count) if tag == "Slides" else "0"
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True).decode("utf-8")


def _update_prs_rels(rels_xml: str, slide_count: int) -> str:
    """ppt/_rels/presentation.xml.rels 에 슬라이드 관계 항목 교체."""
    NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    SLIDE_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    REMOVE_TYPES = {"notesSlide", "notesMaster"}

    root = etree.fromstring(rels_xml.encode("utf-8"))

    # 슬라이드 관계 + 삭제된 notes* 관계 제거
    to_remove = [
        el for el in root
        if el.get("Type") == SLIDE_TYPE
        or any(k in el.get("Type", "") for k in REMOVE_TYPES)
    ]
    for el in to_remove:
        root.remove(el)

    # 남아있는 rId의 최댓값을 구해서 그 이후부터 슬라이드 rId를 부여한다.
    # rId1~N을 그냥 쓰면 slideMaster(rId4), handoutMaster(rId17) 등과 충돌해
    # slideMaster가 사라지면서 PowerPoint가 파일을 열지 못한다.
    existing_nums = []
    for el in root:
        rid = el.get("Id", "")
        if rid.startswith("rId"):
            try:
                existing_nums.append(int(rid[3:]))
            except ValueError:
                pass
    next_id = max(existing_nums, default=0) + 1

    for i in range(1, slide_count + 1):
        rel = etree.SubElement(root, f"{{{NS}}}Relationship")
        rel.set("Id", f"rId{next_id + i - 1}")
        rel.set("Type", SLIDE_TYPE)
        rel.set("Target", f"slides/slide{i}.xml")

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True).decode("utf-8")


# ---------------------------------------------------------------------------
# Stage 2: 차트 데이터 교체
# ---------------------------------------------------------------------------

def _remove_chart_gridlines(chart):
    """차트의 major/minor 격자선 제거."""
    NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    plot_area = chart._element.find(f"{{{NS_C}}}chart/{{{NS_C}}}plotArea")
    if plot_area is None:
        return
    for tag in ["majorGridlines", "minorGridlines"]:
        for elem in plot_area.findall(f".//{{{NS_C}}}{tag}"):
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)


def _patch_chart_series_color(chart, series_idx: int, hex_color: str):
    """차트 시리즈에 단색 채우기 색상 적용."""
    NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    plot_area = chart._element.find(f"{{{NS_C}}}chart/{{{NS_C}}}plotArea")
    if plot_area is None:
        return
    for chart_elem in plot_area:
        ser_list = chart_elem.findall(f"{{{NS_C}}}ser")
        if len(ser_list) > series_idx:
            ser = ser_list[series_idx]
            spPr = ser.find(f"{{{NS_C}}}spPr")
            if spPr is None:
                spPr = etree.SubElement(ser, f"{{{NS_C}}}spPr")
            for old in spPr.findall(f"{{{NS_A}}}solidFill"):
                spPr.remove(old)
            solidFill = etree.SubElement(spPr, f"{{{NS_A}}}solidFill")
            srgbClr = etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr")
            srgbClr.set("val", hex_color.lstrip("#"))


def _post_process_charts(output_path: str, draft: dict, template_path: str, layout_map: dict):
    """
    조립된 PPTX를 python-pptx로 열어 content_chart 슬라이드의 차트 데이터를 draft 값으로 교체.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError:
        print("  [경고] python-pptx 미설치 — 차트 후처리 건너뜀")
        return

    CHART_TYPE_MAP = {
        "bar":            XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked":    XL_CHART_TYPE.COLUMN_STACKED,
        "bar_100":        XL_CHART_TYPE.COLUMN_STACKED_100,
        "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line":           XL_CHART_TYPE.LINE,
        "area":           XL_CHART_TYPE.AREA,
        "pie":            XL_CHART_TYPE.PIE,
        "scatter":        XL_CHART_TYPE.XY_SCATTER_LINES,
    }

    prs = Presentation(output_path)
    slides_data = draft.get("slides", [])
    PAL_CHART = ["#627365", "#2D3734", "#B8CCC4", "#406D92", "#D98F76", "#538184"]

    CHART_LAYOUTS = {"content_chart", "table_chart_combo"}

    for i, slide_data in enumerate(slides_data):
        if i >= len(prs.slides):
            break
        layout = slide_data.get("layout", "")
        if layout not in CHART_LAYOUTS:
            continue

        chart_spec = slide_data.get("chart", {})
        if not chart_spec:
            continue

        slide = prs.slides[i]

        # 1. chart_placeholder (name 기반) 또는 실제 chart 도형 찾기
        chart_shapes = [
            s for s in slide.shapes
            if s.name == "chart_placeholder" or s.shape_type == 3
        ]
        if not chart_shapes:
            continue
        chart_pos = (chart_shapes[0].left, chart_shapes[0].top,
                     chart_shapes[0].width, chart_shapes[0].height)

        # key_points 존재 시 차트 너비 축소 (content_chart only)
        key_points = slide_data.get("key_points", [])
        if key_points and layout == "content_chart":
            chart_pos = (chart_pos[0], chart_pos[1], 7500000, chart_pos[3])

        # 2. placeholder/chart 도형 전부 제거
        sp_tree = slide.shapes._spTree
        for cs in chart_shapes:
            sp_tree.remove(cs._element)

        # 3. 새 차트 데이터 구성
        chart_type_str = chart_spec.get("chart_type", "bar")
        series_data = chart_spec.get("series", chart_spec.get("data", []))
        if not series_data:
            continue

        cd = ChartData()
        multi_series = (
            isinstance(series_data, list)
            and len(series_data) > 0
            and isinstance(series_data[0], dict)
            and "data" in series_data[0]
        )

        series_count = 0
        if multi_series:
            # 다중 시리즈: [{name, data:[{label,value}]}]
            labels = [str(d["label"]) for d in series_data[0]["data"]]
            cd.categories = labels
            for s in series_data:
                cd.add_series(s.get("name", ""), [d["value"] for d in s["data"]])
                series_count += 1
            xl_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.LINE)
        else:
            # 단일 시리즈: [{label, value}]
            labels = [str(d["label"]) for d in series_data]
            values = [d.get("value", 0) for d in series_data]
            cd.categories = labels
            cd.add_series("", values)
            series_count = 1
            xl_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # 4. 새 차트 추가 (같은 위치)
        chart_frame = slide.shapes.add_chart(
            xl_type,
            chart_pos[0], chart_pos[1], chart_pos[2], chart_pos[3],
            cd
        )

        # 5. 스타일 패치
        chart = chart_frame.chart
        chart.has_legend = multi_series or series_count > 1
        _remove_chart_gridlines(chart)
        for idx, color in enumerate(PAL_CHART[:series_count]):
            _patch_chart_series_color(chart, idx, color)

    prs.save(output_path)


# ---------------------------------------------------------------------------
# Stage 3: 표 행 수 조정
# ---------------------------------------------------------------------------

def _post_process_tables(output_path: str, draft: dict, spec: dict | None = None):
    """
    table_slide 슬라이드의 표를 draft 데이터 행 수에 맞게 재구성.
    행 수가 일치하면 텍스트 교체는 replace_placeholders에서 처리된 것으로 간주하고 건너뜀.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from lxml import etree as _etree
    except ImportError:
        print("  [경고] python-pptx 미설치 — 표 후처리 건너뜀")
        return

    prs = Presentation(output_path)
    slides_data = draft.get("slides", [])
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    ts = (spec or {}).get("table_style", {})
    PAL_PRIMARY    = ts.get("header_fill",       "#2D3734")
    PAL_EVEN       = ts.get("row_fill_even",     "#E4E0D4")
    PAL_ODD        = ts.get("row_fill_odd",      "#FFFFFF")
    FONT_TBL_HDR   = int((spec or {}).get("font_sizes", {}).get("table_header_pt", 9))
    FONT_TBL_BODY  = int((spec or {}).get("font_sizes", {}).get("table_body_pt",   9))

    TABLE_LAYOUTS = {"table_slide", "table_chart_combo"}

    for i, slide_data in enumerate(slides_data):
        if i >= len(prs.slides):
            break
        if slide_data.get("layout") not in TABLE_LAYOUTS:
            continue

        table_spec = slide_data.get("table", {})
        headers = table_spec.get("headers", [])
        rows = table_spec.get("rows", [])
        if not headers:
            continue

        slide = prs.slides[i]

        # 기존 TABLE 도형 찾기
        tbl_shape = None
        tbl_pos = None
        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                tbl_shape = shape
                tbl_pos = (shape.left, shape.top, shape.width, shape.height)
                break

        if tbl_shape is None:
            continue

        existing_row_count = len(tbl_shape.table.rows)
        needed_row_count = len(rows) + 1  # +1 for header

        if existing_row_count == needed_row_count:
            # 행 수 일치 → 텍스트 교체는 replace_placeholders에서 처리됨
            continue

        # 행 수 불일치 → 표 제거 후 재생성
        sp_tree = slide.shapes._spTree
        sp_tree.remove(tbl_shape._element)

        n_cols = len(headers)
        n_rows = needed_row_count

        tbl_shape_new = slide.shapes.add_table(
            n_rows, n_cols,
            tbl_pos[0], tbl_pos[1], tbl_pos[2], tbl_pos[3]
        )
        table = tbl_shape_new.table

        # 열 너비: 첫 열 좁게, 나머지 균등
        total_w = tbl_pos[2]
        if n_cols > 1:
            first_w = total_w // (n_cols + 1)
            rest_w = (total_w - first_w) // (n_cols - 1)
        else:
            first_w = rest_w = total_w
        for ci, col in enumerate(table.columns):
            col.width = Emu(first_w if ci == 0 else rest_w)

        # 헤더행
        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = hdr
            tf = cell.text_frame
            for para in tf.paragraphs:
                para.alignment = 2  # CENTER
                for run in para.runs:
                    run.font.name = "Pretendard"
                    run.font.size = Pt(FONT_TBL_HDR)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
            # 헤더 배경색
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = _etree.SubElement(tcPr, f"{{{NS_A}}}solidFill")
            srgbClr = _etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr")
            srgbClr.set("val", PAL_PRIMARY.lstrip("#"))

        # 데이터행
        for ri, row in enumerate(rows):
            row_bg = PAL_EVEN if ri % 2 == 0 else PAL_ODD
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                tf = cell.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.name = "Pretendard Light"
                        run.font.size = Pt(FONT_TBL_BODY)
                        run.font.bold = (ci == 0)
                        run.font.color.rgb = RGBColor(0x23, 0x23, 0x23)
                # 행 배경색
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                solidFill = _etree.SubElement(tcPr, f"{{{NS_A}}}solidFill")
                srgbClr = _etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr")
                srgbClr.set("val", row_bg.lstrip("#"))

    prs.save(output_path)


# ---------------------------------------------------------------------------
# Stage 4: 동적 텍스트 처리
# ---------------------------------------------------------------------------

def _post_process_dynamic_text(output_path: str, draft: dict):
    """
    {{bullets}}, {{key_points}}, {{left_items}}, {{right_items}},
    {{col1_items}} ~ {{col3_items}}, {{takeaways}} 마커를
    draft 데이터의 전체 항목으로 교체.
    항목이 많아도 모두 표시. normAutofit으로 폰트 자동 축소.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("  [경고] python-pptx 미설치 — 동적 텍스트 처리 건너뜀")
        return

    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _build_paragraphs(txBody, items, prefix, font_name, font_size_pt, color_hex, space_after_pt=4):
        """
        기존 단락 제거 후 items로 새 단락 생성.
        템플릿의 첫 번째 단락 pPr/rPr을 상속 → 커스텀 폰트·색상·들여쓰기 자동 유지.
        pPr/rPr 없으면 fallback 파라미터 사용.
        """
        # 템플릿 첫 단락에서 pPr, rPr 추출
        existing_p = txBody.find(f"{{{NS_A}}}p")
        template_pPr = None
        template_rPr = None
        if existing_p is not None:
            template_pPr = existing_p.find(f"{{{NS_A}}}pPr")
            existing_r = existing_p.find(f"{{{NS_A}}}r")
            if existing_r is not None:
                template_rPr = existing_r.find(f"{{{NS_A}}}rPr")

        # 기존 단락 전체 제거
        for p_elem in list(txBody.findall(f"{{{NS_A}}}p")):
            txBody.remove(p_elem)

        for item in items:
            p_elem = etree.SubElement(txBody, f"{{{NS_A}}}p")

            # ── pPr: 템플릿 복사 + lnSpc/spcAft 보장 ──────────────────
            if template_pPr is not None:
                new_pPr = copy.deepcopy(template_pPr)
            else:
                new_pPr = etree.Element(f"{{{NS_A}}}pPr")

            # lnSpc 없으면 150% 추가
            if new_pPr.find(f"{{{NS_A}}}lnSpc") is None:
                lnSpc_el = etree.SubElement(new_pPr, f"{{{NS_A}}}lnSpc")
                spcPct_el = etree.SubElement(lnSpc_el, f"{{{NS_A}}}spcPct")
                spcPct_el.set("val", "150000")

            # spcAft 없으면 추가, 있으면 값 설정
            spcAft_el = new_pPr.find(f"{{{NS_A}}}spcAft")
            if spcAft_el is None:
                spcAft_el = etree.SubElement(new_pPr, f"{{{NS_A}}}spcAft")
            spcPts_el = spcAft_el.find(f"{{{NS_A}}}spcPts")
            if spcPts_el is None:
                spcPts_el = etree.SubElement(spcAft_el, f"{{{NS_A}}}spcPts")
            spcPts_el.set("val", str(int(space_after_pt * 100)))

            p_elem.append(new_pPr)

            # ── rPr: 템플릿 복사 또는 fallback ────────────────────────
            r_elem = etree.SubElement(p_elem, f"{{{NS_A}}}r")
            is_estimated = "[추정]" in item or "(추정)" in item

            if template_rPr is not None:
                new_rPr = copy.deepcopy(template_rPr)
                if is_estimated:
                    for srgb in new_rPr.findall(f".//{{{NS_A}}}srgbClr"):
                        srgb.set("val", "F59E0B")
            else:
                new_rPr = etree.Element(f"{{{NS_A}}}rPr")
                new_rPr.set("lang", "ko-KR")
                new_rPr.set("sz", str(int(font_size_pt * 100)))
                latin = etree.SubElement(new_rPr, f"{{{NS_A}}}latin")
                latin.set("typeface", font_name)
                ea = etree.SubElement(new_rPr, f"{{{NS_A}}}ea")
                ea.set("typeface", font_name)
                item_color = "F59E0B" if is_estimated else color_hex
                solidFill = etree.SubElement(new_rPr, f"{{{NS_A}}}solidFill")
                srgbClr = etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr")
                srgbClr.set("val", item_color.lstrip("#"))

            r_elem.append(new_rPr)
            t_elem = etree.SubElement(r_elem, f"{{{NS_A}}}t")
            t_elem.text = prefix + item

    def _enable_autofit(txBody):
        """텍스트가 상자를 초과할 경우 폰트 자동 축소."""
        bodyPr = txBody.find(f"{{{NS_A}}}bodyPr")
        if bodyPr is None:
            return
        for tag in ["normAutofit", "spAutoFit", "noAutofit"]:
            for el in bodyPr.findall(f"{{{NS_A}}}{tag}"):
                bodyPr.remove(el)
        etree.SubElement(bodyPr, f"{{{NS_A}}}normAutofit")

    def _col_items(slide_data, idx):
        cols = slide_data.get("columns", [])
        return cols[idx].get("items", []) if len(cols) > idx else []

    DYNAMIC_MAP = {
        "content_text": {
            "{{bullets}}": {
                "items_fn": lambda s: s.get("body", []),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 9, "color": "232323",
            },
        },
        "content_chart": {
            "{{key_points}}": {
                "items_fn": lambda s: s.get("key_points", []),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 8, "color": "232323",
            },
        },
        "two_column_compare": {
            "{{left_items}}": {
                "items_fn": lambda s: s.get("column_left", {}).get("items", []),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 9, "color": "232323",
            },
            "{{right_items}}": {
                "items_fn": lambda s: s.get("column_right", {}).get("items", []),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 9, "color": "232323",
            },
        },
        "three_column_summary": {
            "{{col1_items}}": {
                "items_fn": lambda s: _col_items(s, 0),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 9, "color": "232323",
            },
            "{{col2_items}}": {
                "items_fn": lambda s: _col_items(s, 1),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 9, "color": "232323",
            },
            "{{col3_items}}": {
                "items_fn": lambda s: _col_items(s, 2),
                "prefix": "▸  ", "font": "Pretendard Light", "size": 9, "color": "232323",
            },
        },
        "closing_slide": {
            "{{takeaways}}": {
                "items_fn": lambda s: s.get("key_takeaways", []),
                "prefix": "✓  ", "font": "Pretendard Light", "size": 10, "color": "B8CCC4",
            },
        },
    }

    # 자동 탐색용: 마커 이름에서 적절한 prefix 추론
    def _infer_prefix(key: str) -> str:
        if any(x in key for x in ("takeaway", "conclusion", "summary")):
            return "✓  "
        if any(x in key for x in ("key_point", "insight", "highlight")):
            return "●  "
        return "▸  "

    prs = Presentation(output_path)
    slides_data = draft.get("slides", [])
    modified = False

    for i, slide_data in enumerate(slides_data):
        if i >= len(prs.slides):
            break
        layout = slide_data.get("layout", "")
        slide = prs.slides[i]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text

            # ── 1단계: DYNAMIC_MAP 명시 항목 처리 ────────────────────
            if layout in DYNAMIC_MAP:
                for marker, cfg in DYNAMIC_MAP[layout].items():
                    if marker not in full_text:
                        continue
                    items = cfg["items_fn"](slide_data)
                    if not items:
                        continue
                    txBody = shape.text_frame._txBody
                    _enable_autofit(txBody)
                    _build_paragraphs(
                        txBody, items,
                        cfg["prefix"], cfg["font"], cfg["size"], cfg["color"],
                    )
                    modified = True

            # ── 2단계: 자동 탐색 — 잔여 {{marker}}가 리스트면 확장 ──
            remaining = re.findall(r'\{\{[a-z][a-z0-9_.]*\}\}', full_text)
            for marker in remaining:
                key = marker[2:-2]
                val = _auto_resolve(key, slide_data)
                if not isinstance(val, list) or not val:
                    continue
                items = [str(v) for v in val]
                txBody = shape.text_frame._txBody
                _enable_autofit(txBody)
                _build_paragraphs(
                    txBody, items,
                    prefix=_infer_prefix(key),
                    font_name="Pretendard Light",
                    font_size_pt=9,
                    color_hex="232323",
                )
                modified = True

    if modified:
        prs.save(output_path)


# ---------------------------------------------------------------------------
# Stage 5: 존 컴포넌트 삽입
# ---------------------------------------------------------------------------

def _post_process_zones(output_path: str, draft: dict, spec: dict | None = None):
    """
    layout="zone" 또는 zone_config 가 있는 슬라이드에 존 컴포넌트를 삽입한다.
    zone_base 템플릿 슬라이드에서 복사된 빈 슬라이드에 python-pptx 로 표/차트/텍스트를 추가.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from lxml import etree as _etree
    except ImportError:
        print("  [경고] python-pptx 미설치 — 존 후처리 건너뜀")
        return

    _cp  = (spec or {}).get("color_palette", {})
    _ts  = (spec or {}).get("table_style", {})
    _fs  = (spec or {}).get("font_sizes", {})

    PAL_DARK     = _cp.get("color_1", "#2D3734").lstrip("#")
    PAL_PRIMARY  = _ts.get("header_fill", "#627365").lstrip("#")
    PAL_GOLD     = _cp.get("color_4", "#A09567").lstrip("#")
    PAL_TEXT     = RGBColor(0x23, 0x23, 0x23)
    PAL_CHART    = [
        f"#{_cp.get('color_1','2D3734').lstrip('#')}",
        f"#{_cp.get('color_2','B8CCC4').lstrip('#')}",
        f"#{_cp.get('color_3','D4E4DE').lstrip('#')}",
        f"#{_cp.get('color_4','A5948C').lstrip('#')}",
        f"#{_cp.get('color_5','BCB8AF').lstrip('#')}",
        f"#{_cp.get('color_6','E4E0D4').lstrip('#')}",
    ]
    FONT_TBL_HDR  = int(_fs.get("table_header_pt", 9))
    FONT_TBL_BODY = int(_fs.get("table_body_pt",   9))
    TBL_EVEN      = _ts.get("row_fill_even", "#E4E0D4").lstrip("#")
    TBL_ODD       = _ts.get("row_fill_odd",  "#FFFFFF").lstrip("#")
    NS_A         = "http://schemas.openxmlformats.org/drawingml/2006/main"

    CHART_TYPE_MAP = {
        "bar":            XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked":    XL_CHART_TYPE.COLUMN_STACKED,
        "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line":           XL_CHART_TYPE.LINE,
        "area":           XL_CHART_TYPE.AREA,
        "pie":            XL_CHART_TYPE.PIE,
    }

    def _rgb(hex_str: str) -> RGBColor:
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _draw_title_bar(slide, x, y, w, title: str):
        """존 상단 제목 바."""
        from pptx.enum.text import PP_ALIGN
        bar = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(266700))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(PAL_GOLD)
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = "Pretendard SemiBold"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        return 266700 + 76200  # 제목 바 높이 + 여백

    def _draw_bullet_zone(slide, rect: dict, zone: dict):
        """텍스트/불릿 존."""
        from pptx.enum.text import PP_ALIGN
        x, y, w, h = rect["x"], rect["y"], rect["w"], rect["h"]
        title = zone.get("title", "")
        items = zone.get("body", zone.get("items", []))
        if isinstance(items, str):
            items = [items]

        title_offset = 0
        if title:
            title_offset = _draw_title_bar(slide, x, y, w, title)

        content_y = y + title_offset
        content_h = h - title_offset

        txBox = slide.shapes.add_textbox(Emu(x), Emu(content_y), Emu(w), Emu(content_h))
        tf = txBox.text_frame
        tf.word_wrap = True

        # normAutofit
        bodyPr = tf._txBody.find(f"{{{NS_A}}}bodyPr")
        if bodyPr is not None:
            for tag in ["normAutofit", "spAutoFit", "noAutofit"]:
                for el in bodyPr.findall(f"{{{NS_A}}}{tag}"):
                    bodyPr.remove(el)
            _etree.SubElement(bodyPr, f"{{{NS_A}}}normAutofit")

        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = "▸  " + str(item)
            run.font.name = "Pretendard Light"
            run.font.size = Pt(9)
            run.font.color.rgb = PAL_TEXT

    def _draw_table_zone(slide, rect: dict, zone: dict):
        """표 존. 행 높이를 존 높이에 맞게 자동 조정."""
        x, y, w, h = rect["x"], rect["y"], rect["w"], rect["h"]
        title = zone.get("title", "")
        table_data = zone.get("table", {})
        headers = table_data.get("headers", [])
        rows    = table_data.get("rows", [])
        if not headers:
            return

        title_offset = 0
        if title:
            title_offset = _draw_title_bar(slide, x, y, w, title)

        tbl_y  = y + title_offset
        tbl_h  = h - title_offset
        n_rows = len(rows) + 1   # +1 header
        n_cols = len(headers)

        # 행 높이 계산: 존 높이 / 행 수 (최소 228600 EMU = 18pt)
        row_h = max(228600, tbl_h // n_rows)
        actual_h = row_h * n_rows

        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(tbl_y), Emu(w), Emu(actual_h))
        table = tbl_shape.table

        # 열 너비: 첫 열 좁게
        if n_cols > 1:
            first_w = w // (n_cols + 1)
            rest_w  = (w - first_w) // (n_cols - 1)
        else:
            first_w = rest_w = w
        for ci, col in enumerate(table.columns):
            col.width = Emu(first_w if ci == 0 else rest_w)

        NS_A_ETR = "http://schemas.openxmlformats.org/drawingml/2006/main"

        # 헤더행
        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = str(hdr)
            tf = cell.text_frame
            for para in tf.paragraphs:
                para.alignment = 2  # CENTER
                for run in para.runs:
                    run.font.name = "Pretendard"
                    run.font.size = Pt(FONT_TBL_HDR)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            sf = _etree.SubElement(tcPr, f"{{{NS_A_ETR}}}solidFill")
            _etree.SubElement(sf, f"{{{NS_A_ETR}}}srgbClr").set("val", PAL_PRIMARY)

        # 데이터행
        for ri, row in enumerate(rows):
            bg = TBL_EVEN if ri % 2 == 0 else TBL_ODD
            for ci, val in enumerate(row[:n_cols]):
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                tf = cell.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.name = "Pretendard Light"
                        run.font.size = Pt(FONT_TBL_BODY)
                        run.font.bold = (ci == 0)
                        run.font.color.rgb = PAL_TEXT
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                sf = _etree.SubElement(tcPr, f"{{{NS_A_ETR}}}solidFill")
                _etree.SubElement(sf, f"{{{NS_A_ETR}}}srgbClr").set("val", bg)

    def _draw_chart_zone(slide, rect: dict, zone: dict):
        """차트 존. 존 전체를 차트로 채운다."""
        x, y, w, h = rect["x"], rect["y"], rect["w"], rect["h"]
        title = zone.get("title", "")
        chart_spec = zone.get("chart", {})
        if not chart_spec:
            return

        title_offset = 0
        if title:
            title_offset = _draw_title_bar(slide, x, y, w, title)

        chart_y = y + title_offset
        chart_h = h - title_offset

        chart_type_str = chart_spec.get("chart_type", "bar")
        series_data    = chart_spec.get("series", [])
        if not series_data:
            return

        cd = ChartData()
        multi = (
            isinstance(series_data, list)
            and series_data
            and isinstance(series_data[0], dict)
            and "data" in series_data[0]
        )

        if multi:
            cd.categories = [str(d["label"]) for d in series_data[0]["data"]]
            for s in series_data:
                cd.add_series(s.get("name", ""), [d["value"] for d in s["data"]])
            xl_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.LINE)
            n_series = len(series_data)
        else:
            cd.categories = [str(d["label"]) for d in series_data]
            cd.add_series("", [d.get("value", 0) for d in series_data])
            xl_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
            n_series = 1

        chart_frame = slide.shapes.add_chart(xl_type, Emu(x), Emu(chart_y), Emu(w), Emu(chart_h), cd)
        chart = chart_frame.chart
        chart.has_legend = multi or n_series > 1

        # 격자선 제거 + 색상 패치
        _remove_chart_gridlines(chart)
        for idx, color in enumerate(PAL_CHART[:n_series]):
            _patch_chart_series_color(chart, idx, color)

    def _draw_diagram_zone(slide, rect: dict, zone: dict):
        """
        다이어그램 존: 현재는 자유 텍스트박스로 처리.
        향후 python-pptx 도형 조합 확장 가능.
        """
        from pptx.enum.text import PP_ALIGN
        x, y, w, h = rect["x"], rect["y"], rect["w"], rect["h"]
        title = zone.get("title", "")
        content = zone.get("text", zone.get("description", ""))

        title_offset = 0
        if title:
            title_offset = _draw_title_bar(slide, x, y, w, title)

        txBox = slide.shapes.add_textbox(
            Emu(x), Emu(y + title_offset), Emu(w), Emu(h - title_offset)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = content
        run.font.name = "Pretendard Light"
        run.font.size = Pt(9)
        run.font.color.rgb = PAL_TEXT

    COMPONENT_DRAWERS = {
        "bullet":  _draw_bullet_zone,
        "text":    _draw_bullet_zone,
        "table":   _draw_table_zone,
        "chart":   _draw_chart_zone,
        "diagram": _draw_diagram_zone,
    }

    # 캘리브레이션 로드 (분량 검증용, 로깅만)
    cal = load_calibration()

    prs = Presentation(output_path)
    slides_data = draft.get("slides", [])
    modified = False

    for i, slide_data in enumerate(slides_data):
        if i >= len(prs.slides):
            break

        zone_config = slide_data.get("zone_config")
        if not zone_config and slide_data.get("layout") != "zone":
            continue

        # zone_config 미지정이면 zones 수로 추론
        zones = slide_data.get("zones", [])
        if not zones:
            continue
        if not zone_config:
            from zone_layout import suggest_zone_config
            zone_config = suggest_zone_config(zones)

        rects = get_zone_rects(zone_config)
        rect_by_id = {r["id"]: r for r in rects}

        # 분량 초과 경고 (분할은 logic-analyst 단계에서 처리; 여기서는 경고만)
        if should_split_slide(zones, zone_config, cal):
            print(f"  [경고] 슬라이드 {i+1}: 콘텐츠 분량 초과. logic-analyst 단계에서 분할 권장.")

        slide = prs.slides[i]

        # {{source}} 가 없으면 출처 텍스트박스 텍스트를 비움
        source_val = slide_data.get("source", "")
        for shape in slide.shapes:
            if shape.has_text_frame and "Source: {{source}}" in shape.text_frame.text:
                if source_val:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = run.text.replace("{{source}}", source_val)
                else:
                    shape.text_frame.paragraphs[0].runs[0].text = ""
                break

        # 각 존에 컴포넌트 삽입
        for zone in zones:
            zid       = zone.get("id", rects[0]["id"])
            component = zone.get("component", "bullet")
            rect      = rect_by_id.get(zid)

            if rect is None:
                print(f"  [경고] 슬라이드 {i+1}: 존 ID '{zid}' 를 {zone_config} 에서 찾을 수 없음.")
                continue

            drawer = COMPONENT_DRAWERS.get(component)
            if drawer is None:
                print(f"  [경고] 슬라이드 {i+1}: 알 수 없는 컴포넌트 타입 '{component}'.")
                continue

            drawer(slide, rect, zone)

        modified = True

    if modified:
        prs.save(output_path)


# ---------------------------------------------------------------------------
# 동일 레이아웃 복수 슬라이드 — 차트 파일명 재매핑
# ---------------------------------------------------------------------------

def _remap_chart_refs(rels_xml: str, template_zf: zipfile.ZipFile,
                      out_dir: str, src_slide_num: int, dst_slide_num: int) -> str:
    """
    슬라이드 rels XML의 차트 참조를 복사하면서 파일명에 dst_slide_num 접두사 추가.
    ppt/charts/ 에 복사한 차트 파일을 실제로 저장하고 Target 경로를 업데이트.
    """
    NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    try:
        root = etree.fromstring(rels_xml.encode("utf-8"))
    except etree.XMLSyntaxError:
        return rels_xml

    all_names = set(template_zf.namelist())

    for rel in root.findall(f"{{{NS}}}Relationship"):
        target = rel.get("Target", "")
        rel_type = rel.get("Type", "")

        if "chart" not in rel_type.lower() and "chart" not in target.lower():
            continue
        if target.startswith("http"):
            continue

        basename = target.split("/")[-1]  # e.g. chart3.xml
        new_basename = f"s{dst_slide_num}_{basename}"  # e.g. s5_chart3.xml

        # 절대 경로 계산 (ppt/slides/ 기준)
        if target.startswith("../"):
            src_abs = "ppt/" + target[3:]
        elif target.startswith("ppt/"):
            src_abs = target
        else:
            src_abs = "ppt/slides/" + target

        # path normalize
        parts = src_abs.split("/")
        normalized = []
        for p in parts:
            if p == "..":
                if normalized:
                    normalized.pop()
            elif p:
                normalized.append(p)
        src_abs = "/".join(normalized)

        if src_abs not in all_names:
            continue

        new_abs = src_abs.rsplit("/", 1)[0] + "/" + new_basename

        # 차트 XML 복사
        charts_dir_path = os.path.join(out_dir, *new_abs.split("/")[:-1])
        os.makedirs(charts_dir_path, exist_ok=True)
        dst_chart_path = os.path.join(out_dir, *new_abs.split("/"))
        with open(dst_chart_path, "wb") as fp:
            fp.write(template_zf.read(src_abs))

        # 차트 rels 복사
        chart_dir = src_abs.rsplit("/", 1)[0]
        chart_rels_src = chart_dir + "/_rels/" + basename + ".rels"
        if chart_rels_src in all_names:
            chart_rels_dst = chart_dir + "/_rels/" + new_basename + ".rels"
            rels_dir_path = os.path.join(out_dir, *chart_rels_dst.split("/")[:-1])
            os.makedirs(rels_dir_path, exist_ok=True)
            dst_rels_path = os.path.join(out_dir, *chart_rels_dst.split("/"))
            with open(dst_rels_path, "wb") as fp:
                fp.write(template_zf.read(chart_rels_src))

        # Target 업데이트
        rel.set("Target", target.replace(basename, new_basename))

    return etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    ).decode("utf-8")


def _strip_rels_by_type(rels_xml: str, *type_keywords: str) -> str:
    """rels XML에서 특정 Type 키워드를 포함하는 Relationship 요소를 제거한다."""
    NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    try:
        root = etree.fromstring(rels_xml.encode("utf-8"))
    except etree.XMLSyntaxError:
        return rels_xml
    for rel in root.findall(f"{{{NS}}}Relationship"):
        rel_type = rel.get("Type", "").lower()
        if any(kw.lower() in rel_type for kw in type_keywords):
            root.remove(rel)
    return etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    ).decode("utf-8")


def _collect_reachable_parts(out_dir: str) -> set:
    """
    _rels/.rels 에서 시작해 모든 관계를 따라가며 도달 가능한 파트 경로를 반환.
    도달 불가 파트(orphaned)는 Content_Types에서 제거해야 PowerPoint가 정상 열림.
    """
    import posixpath as _pp

    NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
    reachable: set = set()
    queue: list = ["_rels/.rels"]

    def _rels_for(part_path: str) -> str:
        """part_path의 rels 파일 경로 반환."""
        parts = part_path.split("/")
        return "/".join(parts[:-1] + ["_rels", parts[-1] + ".rels"])

    def _resolve(base_dir: str, target: str) -> str:
        if target.startswith("/"):
            return target.lstrip("/")
        return _pp.normpath(_pp.join(base_dir, target))

    visited_rels: set = set()

    while queue:
        rels_path = queue.pop()
        if rels_path in visited_rels:
            continue
        visited_rels.add(rels_path)

        full = os.path.join(out_dir, *rels_path.split("/"))
        if not os.path.exists(full):
            continue

        try:
            root = etree.fromstring(open(full, "rb").read())
        except etree.XMLSyntaxError:
            continue

        # rels 파일이 속한 디렉터리 계산
        # _rels/foo.xml.rels → base_dir is the parent of _rels/
        path_parts = rels_path.split("/")
        rels_idx = path_parts.index("_rels") if "_rels" in path_parts else -1
        if rels_idx >= 0:
            base_dir = "/".join(path_parts[:rels_idx])
        else:
            base_dir = ""

        for rel in root.findall(f"{{{NS_REL}}}Relationship"):
            if rel.get("TargetMode") == "External":
                continue
            target = rel.get("Target", "")
            part_path = _resolve(base_dir, target) if base_dir else _pp.normpath(target)
            reachable.add(part_path)
            # 이 파트의 rels도 큐에 추가
            queue.append(_rels_for(part_path))

    return reachable


def _remove_orphaned_parts(out_dir: str):
    """
    도달 불가 파트(orphaned)를 디스크에서 삭제하고 Content_Types에서 제거.
    PowerPoint는 Content_Types에 등록됐지만 관계 체인에서 참조되지 않는
    파트를 '읽을 수 없는 부분'으로 간주해 복구 메시지를 띄운다.
    """
    reachable = _collect_reachable_parts(out_dir)

    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
    ct_path = os.path.join(out_dir, "[Content_Types].xml")
    root = etree.fromstring(open(ct_path, "rb").read())

    removed_ct = 0
    for el in list(root.findall(f"{{{CT_NS}}}Override")):
        part = el.get("PartName", "").lstrip("/")
        if part and part not in reachable:
            # 디스크에서도 삭제
            disk_path = os.path.join(out_dir, *part.split("/"))
            if os.path.exists(disk_path):
                os.remove(disk_path)
                # rels 파일도 삭제
                p = part.split("/")
                rels_rel = "/".join(p[:-1] + ["_rels", p[-1] + ".rels"])
                rels_disk = os.path.join(out_dir, *rels_rel.split("/"))
                if os.path.exists(rels_disk):
                    os.remove(rels_disk)
            root.remove(el)
            removed_ct += 1

    if removed_ct:
        with open(ct_path, "w", encoding="utf-8") as f:
            f.write(etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            ).decode("utf-8"))
        print(f"  [정리] 고아 파트 {removed_ct}개 제거")


def _update_content_types_with_charts(out_dir: str):
    """
    out_dir을 스캔해 s{N}_chart*.xml 파일을 [Content_Types].xml에 자동 등록.
    """
    ct_path = os.path.join(out_dir, "[Content_Types].xml")
    if not os.path.exists(ct_path):
        return

    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
    CHART_CT = "application/vnd.openxmlformats-officedocument.drawingml.chart+xml"

    with open(ct_path, encoding="utf-8") as f:
        ct_xml = f.read()

    root = etree.fromstring(ct_xml.encode("utf-8"))
    existing_parts = {el.get("PartName") for el in root}

    charts_dir = os.path.join(out_dir, "ppt", "charts")
    if not os.path.isdir(charts_dir):
        return

    added = 0
    for fname in os.listdir(charts_dir):
        if re.match(r"s\d+_chart\d+\.xml$", fname):
            part_name = f"/ppt/charts/{fname}"
            if part_name not in existing_parts:
                override = etree.SubElement(root, f"{{{CT_NS}}}Override")
                override.set("PartName", part_name)
                override.set("ContentType", CHART_CT)
                existing_parts.add(part_name)
                added += 1

    if added > 0:
        with open(ct_path, "w", encoding="utf-8") as f:
            f.write(etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True).decode("utf-8"))


# ---------------------------------------------------------------------------
# PPTX ZIP 기반 조립
# ---------------------------------------------------------------------------

def build_assembled_pptx(draft: dict, template_path: str, template_index: dict, output_path: str):
    """draft 기반으로 template 슬라이드를 복사·편집해 output_path 에 저장."""

    layout_map: dict = template_index.get("layouts", {})  # layout_name -> 1-based slide idx

    # 템플릿 pasteboard 스펙 로드 (폰트·컬러·표 스타일)
    spec = load_spec(template_path)

    slides_data = draft.get("slides", [])

    tmp_dir = tempfile.mkdtemp(prefix="assembler_")
    try:
        # ------------------------------------------------------------------
        # 1. template PPTX 압축 해제
        # ------------------------------------------------------------------
        template_extract_dir = os.path.join(tmp_dir, "template")
        os.makedirs(template_extract_dir, exist_ok=True)
        with zipfile.ZipFile(template_path, "r") as zf:
            zf.extractall(template_extract_dir)
            template_slide_names = _list_slides(zf)

        # ------------------------------------------------------------------
        # 2. 출력 PPTX 디렉토리 구조 준비 (template 복사)
        # ------------------------------------------------------------------
        out_dir = os.path.join(tmp_dir, "output")
        shutil.copytree(template_extract_dir, out_dir)

        # 기존 슬라이드 XML 전부 삭제 (새로 채울 것)
        slides_dir = os.path.join(out_dir, "ppt", "slides")
        slides_rels_dir = os.path.join(out_dir, "ppt", "slides", "_rels")
        for fname in os.listdir(slides_dir):
            if re.match(r"slide\d+\.xml$", fname):
                os.remove(os.path.join(slides_dir, fname))
        if os.path.isdir(slides_rels_dir):
            for fname in os.listdir(slides_rels_dir):
                if re.match(r"slide\d+\.xml\.rels$", fname):
                    os.remove(os.path.join(slides_rels_dir, fname))
        os.makedirs(slides_rels_dir, exist_ok=True)

        # notesSlide 파일 전체 삭제: 슬라이드를 재배치하면 notesSlide의
        # 역방향 참조(slide→notesSlide, notesSlide→slide)가 어긋나
        # PowerPoint가 파일을 열지 못하는 치명적 손상이 발생한다.
        notes_dir = os.path.join(out_dir, "ppt", "notesSlides")
        if os.path.isdir(notes_dir):
            shutil.rmtree(notes_dir)

        # notesMaster도 제거 (notesSlide 없으면 불필요)
        notes_master_dir = os.path.join(out_dir, "ppt", "notesMasters")
        if os.path.isdir(notes_master_dir):
            shutil.rmtree(notes_master_dir)

        # ------------------------------------------------------------------
        # 3. 각 슬라이드 처리
        # ------------------------------------------------------------------
        output_slide_count = 0

        with zipfile.ZipFile(template_path, "r") as template_zf:
            for slide_data in slides_data:
                layout = slide_data.get("layout", "")
                template_slide_number = layout_map.get(layout)

                # zone 슬라이드: zone_config 있거나 layout="zone" → zone_base 템플릿 사용
                if template_slide_number is None and (
                    slide_data.get("zone_config") or layout == "zone"
                ):
                    template_slide_number = layout_map.get("zone_base") or layout_map.get("zone")

                if template_slide_number is None:
                    print(f"  [경고] 레이아웃 '{layout}' 을(를) 인덱스에서 찾을 수 없음. 슬라이드 건너뜀.")
                    continue

                # template 슬라이드 XML 읽기
                template_slide_name = f"ppt/slides/slide{template_slide_number}.xml"
                template_rels_name = f"ppt/slides/_rels/slide{template_slide_number}.xml.rels"

                if template_slide_name not in template_zf.namelist():
                    print(f"  [경고] 템플릿에서 '{template_slide_name}' 를 찾을 수 없음. 슬라이드 건너뜀.")
                    continue

                slide_xml = _read_zip_text(template_zf, template_slide_name)

                # 텍스트 교체
                slide_xml = replace_placeholders(slide_xml, layout, slide_data)

                # 슬라이드 외부 요소 필터
                slide_xml = filter_outside_elements(slide_xml)

                # 출력 슬라이드 번호
                output_slide_count += 1
                out_slide_path = os.path.join(slides_dir, f"slide{output_slide_count}.xml")
                with open(out_slide_path, "w", encoding="utf-8") as f:
                    f.write(slide_xml)

                # rels 파일 처리 (Stage 4: 차트 파일명 재매핑으로 ID 충돌 방지)
                out_rels_path = os.path.join(slides_rels_dir, f"slide{output_slide_count}.xml.rels")
                if template_rels_name in template_zf.namelist():
                    rels_xml = _read_zip_text(template_zf, template_rels_name)
                    # 차트 레이아웃(Stage 2에서 새 차트 삽입)은 _remap_chart_refs를
                    # 건너뛴다: 템플릿 차트 파일을 복사하면 고아 파일이 생겨 PPTX 손상
                    if layout not in {"content_chart", "table_chart_combo"}:
                        rels_xml = _remap_chart_refs(
                            rels_xml, template_zf, out_dir,
                            template_slide_number, output_slide_count
                        )
                    else:
                        # 기존 chart rels 참조만 제거 (파일은 복사하지 않음)
                        rels_xml = _strip_rels_by_type(rels_xml, "chart")
                    # notesSlide 관계 제거: 다른 슬라이드가 같은 notesSlide를
                    # 공유하면 PPTX가 손상되므로 항상 제거한다
                    rels_xml = _strip_rels_by_type(rels_xml, "notesSlide")
                    with open(out_rels_path, "w", encoding="utf-8") as f:
                        f.write(rels_xml)
                else:
                    # 빈 rels 생성
                    with open(out_rels_path, "w", encoding="utf-8") as f:
                        f.write(
                            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
                            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
                        )

        if output_slide_count == 0:
            raise RuntimeError("처리된 슬라이드가 없습니다. draft.json 또는 template_index 를 확인하세요.")

        # ------------------------------------------------------------------
        # 4. [Content_Types].xml 업데이트
        # ------------------------------------------------------------------
        ct_path = os.path.join(out_dir, "[Content_Types].xml")
        with open(ct_path, encoding="utf-8") as f:
            ct_xml = f.read()
        ct_xml = _update_content_types(ct_xml, output_slide_count)
        with open(ct_path, "w", encoding="utf-8") as f:
            f.write(ct_xml)

        # Stage 4: 새로 복사된 s{N}_chart*.xml 파일을 Content_Types에 등록
        _update_content_types_with_charts(out_dir)

        # ------------------------------------------------------------------
        # 5+6. ppt/presentation.xml + ppt/_rels/presentation.xml.rels 업데이트
        # 순서: rels 먼저 → 실제 rId 목록 확정 → presentation.xml에 반영
        # (rId가 기존 관계와 충돌하면 slideMaster 등이 사라져 파일이 열리지 않음)
        # ------------------------------------------------------------------
        prs_rels_path = os.path.join(out_dir, "ppt", "_rels", "presentation.xml.rels")
        slide_rids = [f"rId{i}" for i in range(1, output_slide_count + 1)]  # 기본값
        if os.path.exists(prs_rels_path):
            with open(prs_rels_path, encoding="utf-8") as f:
                prs_rels_xml = f.read()
            prs_rels_xml = _update_prs_rels(prs_rels_xml, output_slide_count)
            with open(prs_rels_path, "w", encoding="utf-8") as f:
                f.write(prs_rels_xml)
            # 실제로 부여된 slide rId 목록 추출 (presentation.xml과 동기화)
            from lxml import etree as _etree2
            NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
            SLIDE_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
            _root = _etree2.fromstring(prs_rels_xml.encode("utf-8"))
            slide_rids = [
                r.get("Id") for r in _root.findall(f"{{{NS_REL}}}Relationship")
                if r.get("Type") == SLIDE_TYPE
            ]

        prs_xml_path = os.path.join(out_dir, "ppt", "presentation.xml")
        with open(prs_xml_path, encoding="utf-8") as f:
            prs_xml = f.read()
        prs_xml = _update_presentation_xml(prs_xml, output_slide_count, slide_rids)
        with open(prs_xml_path, "w", encoding="utf-8") as f:
            f.write(prs_xml)

        # ------------------------------------------------------------------
        # 6b. docProps/app.xml 슬라이드 수 업데이트
        # ------------------------------------------------------------------
        app_xml_path = os.path.join(out_dir, "docProps", "app.xml")
        if os.path.exists(app_xml_path):
            with open(app_xml_path, encoding="utf-8") as f:
                app_xml = f.read()
            app_xml = _update_doc_props(app_xml, output_slide_count)
            with open(app_xml_path, "w", encoding="utf-8") as f:
                f.write(app_xml)

        # ------------------------------------------------------------------
        # 7. 고아 파트 정리 → ZIP 재압축 → output_path
        # ------------------------------------------------------------------
        # 관계 체인에서 참조되지 않는 파트를 Content_Types + 디스크에서 제거.
        # 이것이 없으면 PowerPoint가 "읽을 수 없는 부분 제거" 복구 메시지를 띄움.
        _remove_orphaned_parts(out_dir)

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        tmp_zip = output_path + ".tmp.zip"
        with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
            for root_dir, dirs, files in os.walk(out_dir):
                for file in files:
                    file_abs = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_abs, out_dir).replace(os.sep, "/")
                    zf_out.write(file_abs, arcname)

        if os.path.exists(output_path):
            os.remove(output_path)
        os.rename(tmp_zip, output_path)

        # ------------------------------------------------------------------
        # 8. Stage 2: 차트 데이터 교체 (python-pptx 후처리)
        # ------------------------------------------------------------------
        print("  [Stage 2] 차트 데이터 후처리...")
        _post_process_charts(output_path, draft, template_path, layout_map)

        # ------------------------------------------------------------------
        # 9. Stage 3: 표 행 수 조정 (python-pptx 후처리)
        # ------------------------------------------------------------------
        print("  [Stage 3] 표 행 수 조정 후처리...")
        _post_process_tables(output_path, draft, spec)

        # ------------------------------------------------------------------
        # 10. Stage 4: 동적 텍스트 처리 (bullets, key_points, items, takeaways)
        # ------------------------------------------------------------------
        print("  [Stage 4] 동적 텍스트 처리...")
        _post_process_dynamic_text(output_path, draft)

        # ------------------------------------------------------------------
        # 11. Stage 5: 존 컴포넌트 삽입 (zone_config 슬라이드)
        # ------------------------------------------------------------------
        print("  [Stage 5] 존 컴포넌트 삽입...")
        _post_process_zones(output_path, draft, spec)

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    print(f"  총 {output_slide_count}장 슬라이드 조립 완료.")


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="assembler — draft.json + template.pptx → output.pptx")
    parser.add_argument("--draft",    required=True, help="draft JSON 파일 경로")
    parser.add_argument("--template", required=True, help="component_template PPTX 파일 경로")
    parser.add_argument("--output",   required=True, help="출력 PPTX 파일 경로")
    args = parser.parse_args()

    with open(args.draft, encoding="utf-8") as f:
        draft = json.load(f)

    index_path = args.template.replace(".pptx", "_index.json")
    if not os.path.exists(index_path):
        raise FileNotFoundError(
            f"템플릿 인덱스 파일을 찾을 수 없습니다: {index_path}\n"
            f"먼저 python src/core/build_component_template.py 를 실행하세요."
        )

    with open(index_path, encoding="utf-8") as f:
        template_index = json.load(f)

    print(f"draft   : {args.draft}")
    print(f"template: {args.template}")
    print(f"output  : {args.output}")
    print("조립 시작...")

    build_assembled_pptx(draft, args.template, template_index, args.output)
    print(f"완료: {args.output}")


if __name__ == "__main__":
    main()
