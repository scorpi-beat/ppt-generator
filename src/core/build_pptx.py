# -*- coding: utf-8 -*-
"""
build_pptx.py — SRC Report: 해상풍력발전 동향
슬라이드 크기: 12192000 x 6858000 EMU (960×540pt, 16:9)
"""

import json
import os
import sys
import zipfile
import re
from lxml import etree

# 로컬 모듈 경로 등록 (extract_spec.py)
sys.path.insert(0, os.path.dirname(__file__))
from extract_spec import load_spec

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────
# 슬라이드 크기 (EMU)
# ─────────────────────────────────────────────
SLIDE_W = 12192000
SLIDE_H = 6858000

# ─────────────────────────────────────────────
# 컬러 팔레트 (style_report.json 실측값 기준)
# ─────────────────────────────────────────────
PAL = {
    "primary":       "#627365",   # green-grey (header fill, borders)
    "dark":          "#2D3734",   # deep forest (File2 text/fill)
    "text_main":     "#232323",   # near-black (body text — File1 dominant)
    "gold":          "#A09567",   # warm gold (sub-section labels)
    "accent_warm":   "#D98F76",   # salmon (head-message text, highlight)
    "accent_blue":   "#406D92",   # steel blue (secondary chart, right-col header)
    "teal":          "#538184",   # teal (chart series alt)
    "light_green":   "#B8CCC4",   # mint (dividers, chart series)
    "white":         "#FFFFFF",
    "neutral_light": "#F2F2F2",
    "page_num":      "#808080",
    "chart": ["#627365","#2D3734","#B8CCC4","#406D92","#D98F76","#538184","#B1A160","#8CA1B7"],
}


def _pal_from_spec(spec: dict) -> dict:
    """
    extract_spec 결과(spec)에서 PAL 딕셔너리를 구성한다.
    - spec["color_roles"] : pasteboard 레이블 매핑 (예: "primary" → "#627365")
    - spec["color_palette"]: color_1, color_2, ... 순서 배열 (fallback)
    누락된 키는 기존 PAL 기본값을 유지한다.
    """
    roles   = spec.get("color_roles", {})
    ordered = list(spec.get("color_palette", {}).values())  # x좌표 순 색상 리스트

    def _pick(role_key: str, palette_idx: int, fallback: str) -> str:
        if role_key in roles:
            return roles[role_key]
        if palette_idx < len(ordered):
            return ordered[palette_idx]
        return fallback

    return {
        "primary":       _pick("primary",      0, PAL["primary"]),
        "dark":          _pick("dark",          1, PAL["dark"]),
        "text_main":     _pick("text_main",     2, PAL["text_main"]),
        "gold":          _pick("gold",          3, PAL["gold"]),
        "accent_warm":   _pick("accent_warm",   4, PAL["accent_warm"]),
        "accent_blue":   _pick("accent_blue",   5, PAL["accent_blue"]),
        "teal":          _pick("teal",          6, PAL["teal"]),
        "light_green":   _pick("light_green",   7, PAL["light_green"]),
        "white":         "#FFFFFF",
        "neutral_light": "#F2F2F2",
        "page_num":      "#808080",
        "chart":         ordered if len(ordered) >= 4 else PAL["chart"],
    }


# ─────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def emu(n):
    return Emu(int(n))

def pt(n):
    return Pt(n)

FONT_KO    = "Pretendard"
FONT_LIGHT = "Pretendard Light"
FONT_SEMI  = "Pretendard SemiBold"
FONT_EXTRA = "Pretendard ExtraBold"


def _run_font(run, typeface, size_pt, bold=False, color_hex=None):
    run.font.name = typeface
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color_hex:
        run.font.color.rgb = rgb(color_hex)
    # 동아시아(한국어) 폰트 명시
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', typeface)


def _add_textbox(slide, x, y, w, h):
    txBox = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    txBox.text_frame.word_wrap = True
    return txBox


def _safe_add_textbox(slide, x, y, w, h):
    """슬라이드 외부(음수 좌표) 요소를 차단하는 textbox 추가 helper"""
    if x < 0 or y < 0:
        return None
    return _add_textbox(slide, x, y, w, h)


def set_cell_bg(cell, hex_color):
    """셀 배경색 XML 직접 패치"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for old in tcPr.findall(f'{{{ns}}}solidFill'):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, f'{{{ns}}}solidFill')
    srgbClr   = etree.SubElement(solidFill, f'{{{ns}}}srgbClr')
    srgbClr.set('val', hex_color.lstrip('#'))


def set_cell_text(cell, text, typeface=FONT_LIGHT, size_pt=9, bold=False,
                  color_hex="#232323", align=PP_ALIGN.LEFT):
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            _run_font(run, typeface, size_pt, bold, color_hex)


def set_cell_borders(cell, top=None, bottom=None, left=None, right=None):
    """
    각 면의 선 두께(EMU)와 색상 설정.
    None이면 noFill(테두리 없음).
    top/bottom: (width_emu, hex_color) 튜플
    left/right: None → noFill
    """
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    def _make_line_elem(tag, border_spec):
        existing = tcPr.find(f'{{{ns}}}{tag}')
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, f'{{{ns}}}{tag}')
        if border_spec is None:
            noFill = etree.SubElement(ln, f'{{{ns}}}noFill')  # noqa: F841
        else:
            w_emu, hex_color = border_spec
            ln.set('w', str(int(w_emu)))
            solidFill = etree.SubElement(ln, f'{{{ns}}}solidFill')
            srgbClr   = etree.SubElement(solidFill, f'{{{ns}}}srgbClr')
            srgbClr.set('val', hex_color.lstrip('#'))
        return ln

    _make_line_elem('lnT', top)
    _make_line_elem('lnB', bottom)
    _make_line_elem('lnL', left)
    _make_line_elem('lnR', right)


def _set_cell_padding(cell, marL=9524, marR=9524, marT=9524, marB=0):
    """셀 패딩 설정 (EMU 단위)"""
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set('marL', str(int(marL)))
    tcPr.set('marR', str(int(marR)))
    tcPr.set('marT', str(int(marT)))
    tcPr.set('marB', str(int(marB)))


def add_logo(slide, logo_path):
    """좌하단 로고 삽입 (style_report.json 기준)"""
    if os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            emu(623888), emu(6492899),
            emu(959322), emu(147158)
        )


def add_page_number(slide, number):
    """우하단 페이지 번호 (style_report.json 기준)"""
    tb = _add_textbox(slide, 10424072, 6482324, 1144041, 230832)
    tf = tb.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = str(number)
    _run_font(run, FONT_LIGHT, 9, False, PAL["page_num"])


def add_slide_bg(slide, hex_color):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)


def add_rect(slide, x, y, w, h, fill_hex=None, line_hex=None, line_w_pt=0):
    """직사각형 도형 추가"""
    if x < 0 or y < 0:
        return None
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        emu(x), emu(y), emu(w), emu(h)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
    else:
        shape.fill.background()
    if line_hex and line_w_pt > 0:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_w_pt)
    else:
        shape.line.fill.background()
    return shape


# ─────────────────────────────────────────────
# 헤더 (2단 구조: TABLE 섹션 헤더 + head_message 텍스트박스)
# ─────────────────────────────────────────────

def add_header(slide, section_label, head_message):
    """
    2단 헤더 구조:
    3-A. TABLE 기반 섹션 헤더 바 (y=441327, h=650267)
         - 1행 3열: [섹션번호 | 섹션제목 | 빈칸]
         - 배경 전체 #627365, 테두리 상하 1pt solid black, 좌우 noFill
    3-B. 별도 head_message 텍스트박스 (TABLE 아래)
         - y=1141611, Pretendard SemiBold 14pt, #D98F76
    """
    # ── 3-A. TABLE 섹션 헤더 바 ──────────────────────────────────
    tbl_x  = 614160
    tbl_y  = 441327
    tbl_cx = 10944225
    tbl_cy = 650267

    tbl_shape = slide.shapes.add_table(1, 3, emu(tbl_x), emu(tbl_y), emu(tbl_cx), emu(tbl_cy))
    table = tbl_shape.table

    # 열 너비: 섹션번호(좁), 섹션제목(중간), 우측여백(나머지)
    num_w   = 1200000
    title_w = 5500000
    rest_w  = tbl_cx - num_w - title_w
    table.columns[0].width = emu(num_w)
    table.columns[1].width = emu(title_w)
    table.columns[2].width = emu(rest_w)

    # 행 높이
    table.rows[0].height = emu(tbl_cy)

    # section_label 텍스트에서 섹션번호와 제목을 분리
    # 예) "I. 해상풍력발전 개요" → 번호: "I.", 제목: "해상풍력발전 개요"
    # section_label 전체를 번호셀에 넣고 제목셀은 비워도 되지만
    # draft의 title 필드가 이미 "II-1. 글로벌 시장 동향" 형태로 들어옴
    # → 번호·제목 분리: 첫 공백 이전 토큰이 번호
    parts = section_label.split(' ', 1)
    sec_num   = parts[0] if len(parts) > 0 else ''
    sec_title = parts[1] if len(parts) > 1 else ''

    border_top_bottom = (12700, '#000000')  # 1pt solid black

    for col_idx in range(3):
        cell = table.cell(0, col_idx)
        set_cell_bg(cell, PAL["primary"])
        set_cell_borders(
            cell,
            top=border_top_bottom,
            bottom=border_top_bottom,
            left=None,
            right=None,
        )
        _set_cell_padding(cell, marL=9524, marR=9524, marT=9524, marB=0)

    # 섹션번호 셀
    cell0 = table.cell(0, 0)
    set_cell_text(cell0, sec_num,
                  typeface=FONT_EXTRA, size_pt=19, bold=True,
                  color_hex=PAL["white"], align=PP_ALIGN.LEFT)

    # 섹션제목 셀
    cell1 = table.cell(0, 1)
    set_cell_text(cell1, sec_title,
                  typeface=FONT_EXTRA, size_pt=19, bold=True,
                  color_hex=PAL["white"], align=PP_ALIGN.LEFT)

    # 우측 빈칸 셀은 배경만 설정 (이미 위에서 처리)

    # ── 3-B. head_message 텍스트박스 ─────────────────────────────
    hm_x  = 628044
    hm_y  = 1141611
    hm_cx = 10944225
    hm_cy = 307777

    tb = _add_textbox(slide, hm_x, hm_y, hm_cx, hm_cy)
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = head_message
    _run_font(run, FONT_SEMI, 14, False, PAL["accent_warm"])


def add_source_label(slide, source_text, y_pos=6237289):
    """출처 텍스트 (style_report.json source_label.y_range_start 기준)"""
    tb = _add_textbox(slide, 614160, y_pos, 10944225, 220000)
    tf = tb.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Source: " + source_text
    _run_font(run, FONT_LIGHT, 6, False, PAL["page_num"])


def _patch_chart_series_color(chart, series_idx, hex_color):
    """차트 시리즈 색상 XML 패치"""
    ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    plot_area = chart._element.find(f'{{{ns_c}}}chart/{{{ns_c}}}plotArea')
    if plot_area is None:
        return
    for chart_elem in plot_area:
        ser_list = chart_elem.findall(f'{{{ns_c}}}ser')
        if len(ser_list) > series_idx:
            ser = ser_list[series_idx]
            spPr = ser.find(f'{{{ns_c}}}spPr')
            if spPr is None:
                spPr = etree.SubElement(ser, f'{{{ns_c}}}spPr')
            solidFill = spPr.find(f'{{{ns_a}}}solidFill')
            if solidFill is None:
                solidFill = etree.SubElement(spPr, f'{{{ns_a}}}solidFill')
            srgbClr = solidFill.find(f'{{{ns_a}}}srgbClr')
            if srgbClr is None:
                srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
            srgbClr.set('val', hex_color.lstrip('#'))
            ln = spPr.find(f'{{{ns_a}}}ln')
            if ln is None:
                ln = etree.SubElement(spPr, f'{{{ns_a}}}ln')
            lnSolidFill = ln.find(f'{{{ns_a}}}solidFill')
            if lnSolidFill is None:
                lnSolidFill = etree.SubElement(ln, f'{{{ns_a}}}solidFill')
            lnSrgb = lnSolidFill.find(f'{{{ns_a}}}srgbClr')
            if lnSrgb is None:
                lnSrgb = etree.SubElement(lnSolidFill, f'{{{ns_a}}}srgbClr')
            lnSrgb.set('val', hex_color.lstrip('#'))


def _patch_per_point_colors(chart_obj, colors: list):
    """각 데이터 포인트에 개별 색상(hex 문자열 리스트) 적용."""
    ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    plot_area = chart_obj._element.find(f'{{{ns_c}}}chart/{{{ns_c}}}plotArea')
    if plot_area is None:
        return
    for chart_elem in plot_area:
        for ser in chart_elem.findall(f'{{{ns_c}}}ser'):
            for idx, color in enumerate(colors):
                dPt = etree.SubElement(ser, f'{{{ns_c}}}dPt')
                idx_el = etree.SubElement(dPt, f'{{{ns_c}}}idx')
                idx_el.set('val', str(idx))
                etree.SubElement(dPt, f'{{{ns_c}}}bubble3D').set('val', '0')
                spPr = etree.SubElement(dPt, f'{{{ns_c}}}spPr')
                solidFill = etree.SubElement(spPr, f'{{{ns_a}}}solidFill')
                srgbClr   = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                srgbClr.set('val', color.lstrip('#'))


def _remove_chart_gridlines(chart):
    """차트 격자선 제거"""
    ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    plot_area = chart._element.find(f'{{{ns_c}}}chart/{{{ns_c}}}plotArea')
    if plot_area is None:
        return
    for tag in ['majorGridlines', 'minorGridlines']:
        for elem in plot_area.findall(f'.//{{{ns_c}}}{tag}'):
            elem.getparent().remove(elem)


def _set_chart_font_sizes(chart, axis_pt=7, legend_pt=7):
    """차트 축 레이블 및 범례 폰트 크기 설정 (단위: pt)"""
    ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    chart_root = chart._element

    def _apply_txPr(parent, size_pt):
        txPr = parent.find(f'{{{ns_c}}}txPr')
        if txPr is None:
            txPr = etree.SubElement(parent, f'{{{ns_c}}}txPr')
        bodyPr = txPr.find(f'{{{ns_a}}}bodyPr')
        if bodyPr is None:
            etree.SubElement(txPr, f'{{{ns_a}}}bodyPr')
        lstStyle = txPr.find(f'{{{ns_a}}}lstStyle')
        if lstStyle is None:
            etree.SubElement(txPr, f'{{{ns_a}}}lstStyle')
        p = txPr.find(f'{{{ns_a}}}p')
        if p is None:
            p = etree.SubElement(txPr, f'{{{ns_a}}}p')
        pPr = p.find(f'{{{ns_a}}}pPr')
        if pPr is None:
            pPr = etree.SubElement(p, f'{{{ns_a}}}pPr')
        defRPr = pPr.find(f'{{{ns_a}}}defRPr')
        if defRPr is None:
            defRPr = etree.SubElement(pPr, f'{{{ns_a}}}defRPr')
        defRPr.set('sz', str(int(size_pt * 100)))
        defRPr.set('b', '0')

    # 축 (catAx, valAx, serAx)
    for ax_tag in ['catAx', 'valAx', 'serAx']:
        for ax in chart_root.findall(f'.//{{{ns_c}}}{ax_tag}'):
            _apply_txPr(ax, axis_pt)

    # 범례
    legend = chart_root.find(f'.//{{{ns_c}}}legend')
    if legend is not None:
        _apply_txPr(legend, legend_pt)


def _add_key_points_panel(slide, key_points, x_kp, y_start, cx_kp, cy):
    """key_points 우측 패널"""
    if not key_points:
        return
    add_rect(slide, x_kp - 80000, y_start, 12700, cy, fill_hex=PAL["primary"])
    tb = _safe_add_textbox(slide, x_kp, y_start + 100000, cx_kp, cy - 100000)
    if tb is None:
        return
    tf = tb.text_frame
    tf.margin_top = emu(50000)
    first = True
    for kp in key_points:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(5)
        run = para.add_run()
        run.text = "▸  " + kp
        _run_font(run, FONT_LIGHT, 8, False, PAL["text_main"])


# ─────────────────────────────────────────────
# 슬라이드 빌더
# ─────────────────────────────────────────────

def build_title_slide(prs, slide_data, logo_path):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["dark"])

    # 상단 장식선
    add_rect(slide, 0, 1200000, SLIDE_W, 8000, fill_hex=PAL["primary"])

    # 제목
    tb = _add_textbox(slide, 1200000, 2000000, 9800000, 1200000)
    tf = tb.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = slide_data.get('title', '')
    _run_font(run, FONT_SEMI, 32, True, PAL["white"])

    # 구분선
    add_rect(slide, 1200000, 3300000, 4000000, 6000, fill_hex=PAL["primary"])

    # 부제목
    tb2 = _add_textbox(slide, 1200000, 3400000, 9800000, 800000)
    tf2 = tb2.text_frame
    para2 = tf2.paragraphs[0]
    para2.alignment = PP_ALIGN.LEFT
    run2 = para2.add_run()
    run2.text = slide_data.get('subtitle', '')
    _run_font(run2, FONT_LIGHT, 13, False, PAL["light_green"])

    # 날짜
    tb3 = _add_textbox(slide, 1200000, 4600000, 3000000, 300000)
    tf3 = tb3.text_frame
    para3 = tf3.paragraphs[0]
    para3.alignment = PP_ALIGN.LEFT
    run3 = para3.add_run()
    run3.text = slide_data.get('date', '')
    _run_font(run3, FONT_LIGHT, 9, False, PAL["light_green"])

    add_logo(slide, logo_path)
    return slide


def build_content_text(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    body_items = slide_data.get('body', [])
    # 본문 시작 y: 1454151 (style_report.json 기준)
    tb = _add_textbox(slide, 628044, 1454151, 10973754, 5050000)
    tf = tb.text_frame
    tf.margin_top = emu(50000)
    tf.margin_left = emu(50000)

    first = True
    for item in body_items:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(6)
        pPr = para._p.get_or_add_pPr()
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is None:
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', '130000')

        run = para.add_run()
        run.text = "▸  " + item
        if "[추정]" in item or "(추정)" in item:
            _run_font(run, FONT_LIGHT, 9, False, "#F59E0B")
        else:
            _run_font(run, FONT_LIGHT, 9, False, PAL["text_main"])

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


def build_content_chart(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    chart_spec = slide_data.get('chart', {})
    chart_type = chart_spec.get('chart_type', 'bar')
    key_points = slide_data.get('key_points', [])

    # 차트 영역 좌표 (style_report.json 기준)
    if key_points:
        x_chart  = 614160
        y_chart  = 1773238
        cx_chart = 7500000
        cy_chart = 4462510
    else:
        x_chart  = 614160
        y_chart  = 1773238
        cx_chart = 10944225
        cy_chart = 4462510

    series_data = chart_spec.get('series', chart_spec.get('data', []))
    cd = ChartData()
    labels = [str(d['label']) for d in series_data]
    cd.categories = labels

    if isinstance(series_data[0], dict) and 'data' in series_data[0]:
        for s in series_data:
            cd.add_series(s.get('name',''), [d['value'] for d in s['data']])
        actual_type = XL_CHART_TYPE.LINE
    else:
        values = [d.get('value', 0) for d in series_data]
        cd.add_series('', values)
        actual_type = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_type == 'bar' else XL_CHART_TYPE.LINE

    chart_obj = slide.shapes.add_chart(
        actual_type,
        emu(x_chart), emu(y_chart), emu(cx_chart), emu(cy_chart), cd
    ).chart
    chart_obj.has_legend = len(cd.series) > 1
    _remove_chart_gridlines(chart_obj)
    _set_chart_font_sizes(chart_obj)
    _patch_chart_series_color(chart_obj, 0, PAL["primary"])
    if len(cd.series) > 1:
        _patch_chart_series_color(chart_obj, 1, PAL["accent_blue"])

    x_kp  = x_chart + cx_chart + 80000
    cx_kp = SLIDE_W - x_kp - 200000
    _add_key_points_panel(slide, key_points, x_kp, y_chart, cx_kp, cy_chart)

    source = chart_spec.get('source', '')
    if source:
        add_source_label(slide, source)

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


def _build_col_content(slide, col_data, x, y, w, h, title_color=None):
    if not col_data:
        return
    tc = title_color or PAL["primary"]
    title_bar = add_rect(slide, x, y, w, 400000, fill_hex=tc)
    if title_bar is None:
        return
    tf = title_bar.text_frame
    tf.word_wrap = False
    tf.margin_left = emu(76200)
    tf.margin_top  = emu(80000)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = col_data.get('title', '')
    _run_font(run, FONT_SEMI, 11, True, PAL["white"])

    items = col_data.get('items', [])
    tb = _safe_add_textbox(slide, x + 60000, y + 480000, w - 120000, h - 480000)
    if tb is None:
        return
    tf2 = tb.text_frame
    tf2.margin_top = emu(30000)
    first = True
    for item in items:
        if first:
            para2 = tf2.paragraphs[0]
            first = False
        else:
            para2 = tf2.add_paragraph()
        para2.alignment = PP_ALIGN.LEFT
        para2.space_before = Pt(5)
        run2 = para2.add_run()
        run2.text = "▸  " + item
        _run_font(run2, FONT_LIGHT, 9, False, PAL["text_main"])


def build_two_column_compare(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    # 본문 시작 y: 1454151
    y_col   = 1454151
    h_col   = SLIDE_H - y_col - 350000
    x_left  = 614160
    w_left  = 5391354
    x_right = 6277940
    w_right = 5290172

    _build_col_content(slide, slide_data.get('column_left'),
                       x_left, y_col, w_left, h_col, PAL["primary"])
    _build_col_content(slide, slide_data.get('column_right'),
                       x_right, y_col, w_right, h_col, PAL["accent_blue"])

    # 구분선
    add_rect(slide, x_left + w_left + 50000, y_col, 12700, h_col,
             fill_hex=PAL["light_green"])

    source = slide_data.get('source', '')
    if source:
        add_source_label(slide, source)

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


def build_three_column_summary(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    columns = slide_data.get('columns', [])
    n = len(columns)
    if n == 0:
        add_logo(slide, logo_path)
        add_page_number(slide, slide_number)
        return slide

    total_w = 10944225
    col_w   = total_w // n
    x_start = 614160
    # 본문 시작 y: 1454151
    y_col   = 1454151
    h_col   = SLIDE_H - y_col - 350000

    for i, col in enumerate(columns):
        x = x_start + i * col_w
        # 상단 컬러 선
        add_rect(slide, x + 40000, y_col, col_w - 80000, 15000, fill_hex=PAL["primary"])
        # 열 제목
        tb_title = _add_textbox(slide, x + 40000, y_col + 50000, col_w - 80000, 400000)
        tf_t = tb_title.text_frame
        para_t = tf_t.paragraphs[0]
        para_t.alignment = PP_ALIGN.LEFT
        run_t = para_t.add_run()
        run_t.text = col.get('title', '')
        _run_font(run_t, FONT_SEMI, 11, True, PAL["primary"])
        # 항목
        items = col.get('items', [])
        tb_body = _add_textbox(slide, x + 40000, y_col + 480000,
                               col_w - 80000, h_col - 500000)
        tf_b = tb_body.text_frame
        tf_b.margin_top = emu(30000)
        first = True
        for item in items:
            if first:
                para_b = tf_b.paragraphs[0]
                first = False
            else:
                para_b = tf_b.add_paragraph()
            para_b.alignment = PP_ALIGN.LEFT
            para_b.space_before = Pt(5)
            run_b = para_b.add_run()
            run_b.text = "▸  " + item
            _run_font(run_b, FONT_LIGHT, 9, False, PAL["text_main"])
        # 열 사이 구분선
        if i < n - 1:
            add_rect(slide, x + col_w - 6350, y_col, 12700, h_col,
                     fill_hex=PAL["light_green"])

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


def build_table_slide(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    table_data = slide_data.get('table', {})
    headers = table_data.get('headers', [])
    rows    = table_data.get('rows', [])

    if not headers:
        add_logo(slide, logo_path)
        add_page_number(slide, slide_number)
        return slide

    n_cols = len(headers)
    n_rows = len(rows) + 1
    key_points = slide_data.get('key_points', [])

    if key_points:
        table_cx = 7800000
        x_kp     = 8000000
        cx_kp    = 3900000
    else:
        table_cx = 10944225
        x_kp     = None
        cx_kp    = None

    # 표 시작 y: 1773238 (style_report.json table_area.y 기준)
    table_y      = 1773238
    SLIDE_BOTTOM = 6580000   # 슬라이드 하단 여백 위
    HDR_ROW_H    = 480000    # 헤더 행 높이 (EMU)
    DATA_ROW_H   = 420000    # 데이터 행 높이 (EMU)
    max_table_cy = SLIDE_BOTTOM - table_y
    calc_cy      = HDR_ROW_H + DATA_ROW_H * len(rows)
    table_cy     = min(calc_cy, max_table_cy)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        emu(614160), emu(table_y),
        emu(table_cx), emu(table_cy)
    )
    table = tbl_shape.table

    # 열 너비: 첫 열 약간 좁게
    if n_cols > 1:
        first_w = table_cx // (n_cols + 1)
        rest_w  = (table_cx - first_w) // (n_cols - 1)
    else:
        first_w = table_cx
        rest_w  = table_cx
    for i, col in enumerate(table.columns):
        col.width = emu(first_w if i == 0 else rest_w)

    # 행 높이 명시
    table.rows[0].height = emu(HDR_ROW_H)
    for r_i in range(len(rows)):
        table.rows[r_i + 1].height = emu(DATA_ROW_H)

    # 헤더행
    for c_i, hdr in enumerate(headers):
        cell = table.cell(0, c_i)
        set_cell_bg(cell, PAL["primary"])
        set_cell_borders(
            cell,
            top=(12700, '#000000'),
            bottom=(12700, '#000000'),
            left=None,
            right=None,
        )
        _set_cell_padding(cell, marL=9524, marR=9524, marT=9524, marB=0)
        set_cell_text(cell, hdr, typeface=FONT_KO, size_pt=9,
                      bold=True, color_hex=PAL["white"], align=PP_ALIGN.CENTER)

    # 데이터행
    for r_i, row in enumerate(rows):
        row_bg = PAL["neutral_light"] if r_i % 2 == 0 else PAL["white"]
        for c_i, val in enumerate(row):
            cell = table.cell(r_i + 1, c_i)
            set_cell_bg(cell, row_bg)
            set_cell_borders(
                cell,
                top=(12700, '#000000'),
                bottom=(12700, '#000000'),
                left=None,
                right=None,
            )
            _set_cell_padding(cell, marL=9524, marR=9524, marT=9524, marB=0)
            if c_i == 0:
                set_cell_text(cell, str(val), typeface=FONT_SEMI, size_pt=9,
                              bold=True, color_hex=PAL["primary"])
            else:
                set_cell_text(cell, str(val), typeface=FONT_LIGHT, size_pt=9,
                              color_hex=PAL["text_main"])

    if key_points and x_kp:
        _add_key_points_panel(slide, key_points, x_kp, table_y, cx_kp, table_cy)

    source = table_data.get('source', '')
    if source:
        add_source_label(slide, source)

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


def build_roadmap_timeline(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    timeline = slide_data.get('timeline', [])
    n = len(timeline)
    if n == 0:
        add_logo(slide, logo_path)
        add_page_number(slide, slide_number)
        return slide

    line_x_start = 800000
    line_x_end   = 11400000
    line_y       = 3500000
    line_len     = line_x_end - line_x_start
    step_w       = line_len // (n - 1) if n > 1 else line_len

    # 기준선
    add_rect(slide, line_x_start, line_y - 6350, line_len, 12700,
             fill_hex=PAL["primary"])

    for i, item in enumerate(timeline):
        x_center = line_x_start + i * step_w
        node_r   = 200000

        # 노드
        node = slide.shapes.add_shape(
            9,  # OVAL
            emu(x_center - node_r), emu(line_y - node_r),
            emu(node_r * 2), emu(node_r * 2)
        )
        is_first = (i == 0)
        is_last  = (i == n - 1)
        node_color = PAL["primary"] if (is_first or is_last) else PAL["accent_warm"]
        node.fill.solid()
        node.fill.fore_color.rgb = rgb(node_color)
        node.line.color.rgb = rgb(PAL["white"])
        node.line.width = Pt(1.5)

        col_w = min(step_w - 50000, 1700000)
        col_x = max(x_center - col_w // 2, line_x_start)

        if i % 2 == 0:
            # 위 배치
            desc_y   = 1700000
            title_y  = 2200000
            period_y = 2750000
        else:
            # 아래 배치
            period_y = 3800000
            title_y  = 4100000
            desc_y   = 4500000

        # 기간
        tb_p = _add_textbox(slide, col_x, period_y, col_w, 280000)
        para_p = tb_p.text_frame.paragraphs[0]
        para_p.alignment = PP_ALIGN.CENTER
        run_p = para_p.add_run()
        run_p.text = item.get('period', '')
        _run_font(run_p, FONT_SEMI, 9, True, PAL["primary"])

        # 제목
        tb_t = _add_textbox(slide, col_x, title_y, col_w, 380000)
        tb_t.text_frame.word_wrap = True
        para_t = tb_t.text_frame.paragraphs[0]
        para_t.alignment = PP_ALIGN.CENTER
        run_t = para_t.add_run()
        run_t.text = item.get('title', '')
        _run_font(run_t, FONT_SEMI, 8, True, PAL["text_main"])

        # 설명
        tb_d = _add_textbox(slide, col_x, desc_y, col_w, 650000)
        para_d = tb_d.text_frame.paragraphs[0]
        para_d.alignment = PP_ALIGN.CENTER
        run_d = para_d.add_run()
        run_d.text = item.get('description', '')
        _run_font(run_d, FONT_LIGHT, 7, False, PAL["text_main"])

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


def build_closing_slide(prs, slide_data, logo_path):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["dark"])

    # 장식선
    add_rect(slide, 0, 1800000, SLIDE_W, 8000, fill_hex=PAL["primary"])

    # 제목
    tb_title = _add_textbox(slide, 1200000, 1900000, 9800000, 1000000)
    tf_t = tb_title.text_frame
    para_t = tf_t.paragraphs[0]
    para_t.alignment = PP_ALIGN.LEFT
    run_t = para_t.add_run()
    run_t.text = slide_data.get('title', '')
    _run_font(run_t, FONT_SEMI, 22, True, PAL["white"])

    closing_msg = slide_data.get('closing_message', '')
    if closing_msg:
        tb_msg = _add_textbox(slide, 1200000, 3000000, 9800000, 700000)
        tf_m = tb_msg.text_frame
        para_m = tf_m.paragraphs[0]
        para_m.alignment = PP_ALIGN.LEFT
        run_m = para_m.add_run()
        run_m.text = closing_msg
        _run_font(run_m, FONT_LIGHT, 11, False, PAL["light_green"])

    takeaways = slide_data.get('key_takeaways', [])
    tb_tk = _add_textbox(slide, 1200000, 3800000, 9800000, 1800000)
    tf_tk = tb_tk.text_frame
    first = True
    for tw in takeaways:
        if first:
            para_tk = tf_tk.paragraphs[0]
            first = False
        else:
            para_tk = tf_tk.add_paragraph()
        para_tk.alignment = PP_ALIGN.LEFT
        para_tk.space_before = Pt(5)
        run_tk = para_tk.add_run()
        run_tk.text = "✓  " + tw
        _run_font(run_tk, FONT_LIGHT, 10, False, PAL["light_green"])

    disclaimer = slide_data.get('disclaimer', '')
    if disclaimer:
        tb_disc = _add_textbox(slide, 1200000, 5900000, 9800000, 400000)
        tf_disc = tb_disc.text_frame
        para_disc = tf_disc.paragraphs[0]
        para_disc.alignment = PP_ALIGN.LEFT
        run_disc = para_disc.add_run()
        run_disc.text = disclaimer
        _run_font(run_disc, FONT_LIGHT, 6, False, PAL["page_num"])

    add_logo(slide, logo_path)
    return slide


# ─────────────────────────────────────────────
# 슬라이드 4 — bar chart 카테고리별 색상
# ─────────────────────────────────────────────
def build_slide4_bar_chart(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    chart_spec = slide_data.get('chart', {})
    series_data = chart_spec.get('series', chart_spec.get('data', []))
    labels = [str(d['label']) for d in series_data]
    values = [d.get('value', 0) for d in series_data]

    cd = ChartData()
    cd.categories = labels
    cd.add_series('', values)

    key_points = slide_data.get('key_points', [])
    if key_points:
        cx_chart = 7500000
    else:
        cx_chart = 10944225

    x_chart  = 614160
    y_chart  = 1773238
    cy_chart = 4462510

    chart_obj = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        emu(x_chart), emu(y_chart), emu(cx_chart), emu(cy_chart), cd
    ).chart
    chart_obj.has_legend = False
    _remove_chart_gridlines(chart_obj)
    _set_chart_font_sizes(chart_obj)

    # 개별 막대 색상 — 실적/추정/전망
    cat_map = {
        '실적': '627365', '추정': 'B8CCC4', '전망': '406D92'
    }
    point_colors = [cat_map.get(sd.get('category', 'default'), '627365') for sd in series_data]
    _patch_per_point_colors(chart_obj, point_colors)

    x_kp  = x_chart + cx_chart + 80000
    cx_kp = SLIDE_W - x_kp - 200000
    _add_key_points_panel(slide, key_points, x_kp, y_chart, cx_kp, cy_chart)

    source = chart_spec.get('source', '')
    if source:
        add_source_label(slide, source)

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


# ─────────────────────────────────────────────
# 슬라이드 7 — line chart (2-series)
# ─────────────────────────────────────────────
def build_slide7_line_chart(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    chart_spec = slide_data.get('chart', {})
    series_list = chart_spec.get('series', [])
    cd = ChartData()

    if series_list and isinstance(series_list[0], dict) and 'data' in series_list[0]:
        labels = [str(d['label']) for d in series_list[0]['data']]
        cd.categories = labels
        for s in series_list:
            vals = [d['value'] for d in s['data']]
            cd.add_series(s.get('name', ''), vals)
    else:
        data_list = chart_spec.get('data', [])
        labels = [str(d['label']) for d in data_list]
        cd.categories = labels
        cd.add_series('단지 용량 (MW)', [d.get('site_capacity_mw', 0) for d in data_list])
        cd.add_series('터빈 용량 (MW)', [d.get('turbine_mw', 0) for d in data_list])

    key_points = slide_data.get('key_points', [])
    if key_points:
        cx_chart = 7500000
    else:
        cx_chart = 10944225

    x_chart  = 614160
    y_chart  = 1773238
    cy_chart = 4462510

    chart_obj = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        emu(x_chart), emu(y_chart), emu(cx_chart), emu(cy_chart), cd
    ).chart
    chart_obj.has_legend = True
    _remove_chart_gridlines(chart_obj)
    _set_chart_font_sizes(chart_obj)
    _patch_chart_series_color(chart_obj, 0, PAL["primary"])
    _patch_chart_series_color(chart_obj, 1, PAL["accent_warm"])

    x_kp  = x_chart + cx_chart + 80000
    cx_kp = SLIDE_W - x_kp - 200000
    _add_key_points_panel(slide, key_points, x_kp, y_chart, cx_kp, cy_chart)

    source = chart_spec.get('source', '')
    if source:
        add_source_label(slide, source)

    estimated = slide_data.get('estimated_values', [])
    if estimated:
        tb_est = _add_textbox(slide, x_chart, 6050000, cx_chart, 180000)
        para_est = tb_est.text_frame.paragraphs[0]
        run_est = para_est.add_run()
        run_est.text = "* " + estimated[0]
        _run_font(run_est, FONT_LIGHT, 6, False, PAL["page_num"])

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


# ─────────────────────────────────────────────
# 슬라이드 18 — funnel bar chart
# ─────────────────────────────────────────────
def build_slide18_funnel_chart(prs, slide_data, logo_path, slide_number):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_bg(slide, PAL["white"])

    add_header(slide, slide_data.get('title', ''), slide_data.get('head_message', ''))

    chart_spec = slide_data.get('chart', {})
    series_data = chart_spec.get('series', chart_spec.get('data', []))
    labels = [str(d['label']) for d in series_data]
    values = [d.get('value', 0) for d in series_data]

    cd = ChartData()
    cd.categories = labels
    cd.add_series('', values)

    key_points = slide_data.get('key_points', [])
    if key_points:
        cx_chart = 7500000
    else:
        cx_chart = 10944225

    x_chart  = 614160
    y_chart  = 1773238
    cy_chart = 4462510

    chart_obj = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        emu(x_chart), emu(y_chart), emu(cx_chart), emu(cy_chart), cd
    ).chart
    chart_obj.has_legend = False
    _remove_chart_gridlines(chart_obj)
    _set_chart_font_sizes(chart_obj)

    # 단계별 개별 색상 (심각성 그라데이션)
    funnel_colors = ['627365', '4A5E4D', '7A9B8A', 'B8CCC4', 'D98F76']
    _patch_per_point_colors(chart_obj, funnel_colors[:len(series_data)])

    x_kp  = x_chart + cx_chart + 80000
    cx_kp = SLIDE_W - x_kp - 200000
    _add_key_points_panel(slide, key_points, x_kp, y_chart, cx_kp, cy_chart)

    source = chart_spec.get('source', '')
    if source:
        add_source_label(slide, source)

    add_logo(slide, logo_path)
    add_page_number(slide, slide_number)
    return slide


# ─────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────
def main():
    import argparse
    parser = argparse.ArgumentParser(description="해상풍력 리포트 PPTX 빌드")
    parser.add_argument("--draft",    default=None, help="draft JSON 경로")
    parser.add_argument("--output",   default=None, help="출력 PPTX 경로")
    parser.add_argument("--template", default=None,
                        help="스펙 추출 기준 PPTX (pasteboard 팔레트 소스)")
    args = parser.parse_args()

    base_dir    = os.path.dirname(os.path.dirname(os.path.dirname(
                    os.path.abspath(__file__))))
    draft_path  = args.draft  or os.path.join(base_dir, 'outputs', 'draft_report_offshore_wind.json')
    output_path = args.output or os.path.join(base_dir, 'outputs', 'report_offshore_wind.pptx')
    logo_path   = os.path.join(base_dir, 'outputs', 'logo_src.png')

    # 템플릿 PPTX에서 컬러 팔레트 로드 — 지정 없으면 component_template 사용
    template_path = args.template or os.path.join(
        base_dir, 'outputs', 'component_template_report.pptx'
    )
    if os.path.exists(template_path):
        spec = load_spec(template_path)
        PAL.update(_pal_from_spec(spec))
        print(f"  [팔레트] {template_path} 에서 로드")
    else:
        print("  [팔레트] 템플릿 없음 — 기본값 사용")

    with open(draft_path, encoding='utf-8') as f:
        draft = json.load(f)

    slides_data    = draft['slides']
    expected_count = len(slides_data)
    print(f"슬라이드 데이터 로드: {expected_count}장")

    prs = Presentation()
    prs.slide_width  = emu(SLIDE_W)
    prs.slide_height = emu(SLIDE_H)

    for sd in slides_data:
        n   = sd['slide_number']
        lay = sd['layout']
        print(f"  슬라이드 {n}: {lay}")

        if lay == 'title_slide':
            build_title_slide(prs, sd, logo_path)
        elif lay == 'content_text':
            build_content_text(prs, sd, logo_path, n)
        elif lay == 'content_chart':
            if n == 4:
                build_slide4_bar_chart(prs, sd, logo_path, n)
            elif n == 7:
                build_slide7_line_chart(prs, sd, logo_path, n)
            elif n == 18:
                build_slide18_funnel_chart(prs, sd, logo_path, n)
            else:
                build_content_chart(prs, sd, logo_path, n)
        elif lay == 'two_column_compare':
            build_two_column_compare(prs, sd, logo_path, n)
        elif lay == 'three_column_summary':
            build_three_column_summary(prs, sd, logo_path, n)
        elif lay == 'table_slide':
            build_table_slide(prs, sd, logo_path, n)
        elif lay == 'roadmap_timeline':
            build_roadmap_timeline(prs, sd, logo_path, n)
        elif lay == 'closing_slide':
            build_closing_slide(prs, sd, logo_path)
        else:
            print(f"    [WARN] 알 수 없는 레이아웃 '{lay}', content_text 대체")
            build_content_text(prs, sd, logo_path, n)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"\n저장 완료: {output_path}")

    actual = len([f for f in zipfile.ZipFile(output_path).namelist()
                  if re.match(r'ppt/slides/slide\d+\.xml', f)])
    assert actual == expected_count, \
        f"슬라이드 수 불일치: 예상 {expected_count} vs 실제 {actual}"
    print(f"슬라이드 수 검증 통과: {actual}장")
    print("완료.")


if __name__ == '__main__':
    main()
