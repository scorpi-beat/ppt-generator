"""
build_component_template.py — 12개 레이아웃 타입의 디자인 원본 PPTX 생성

실행:
  python src/core/build_component_template.py

출력:
  outputs/component_template_report.pptx
  outputs/component_template_report_index.json
"""

import json
import os

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# 슬라이드 크기
# ---------------------------------------------------------------------------
SLIDE_W = 12192000
SLIDE_H = 6858000

# ---------------------------------------------------------------------------
# 팔레트
# ---------------------------------------------------------------------------
PAL = {
    "primary":       "#627365",
    "dark":          "#2D3734",
    "text_main":     "#232323",
    "gold":          "#A09567",
    "accent_warm":   "#D98F76",
    "accent_blue":   "#406D92",
    "teal":          "#538184",
    "light_green":   "#B8CCC4",
    "white":         "#FFFFFF",
    "neutral_light": "#F2F2F2",
    "page_num":      "#808080",
    "chart_colors":  ["#627365", "#2D3734", "#B8CCC4", "#406D92", "#D98F76", "#538184"],
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# 헬퍼 함수
# ---------------------------------------------------------------------------

def set_cell_border(cell, top=None, bottom=None, left=None, right=None):
    """셀에 테두리 설정. 각 인자는 (width_emu, hex_color) 또는 None(noFill)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    def _make_ln(w_emu, hex_color):
        ln = etree.SubElement(tcPr, qn("a:ln"))
        ln.set("w", str(int(w_emu)))
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", hex_color.lstrip("#"))
        return ln

    def _make_ln_none():
        ln = etree.SubElement(tcPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
        return ln

    # 기존 lnXxx 제거
    for tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)

    def _add_border(tag, spec):
        if spec is None:
            ln = _make_ln_none()
        else:
            ln = _make_ln(spec[0], spec[1])
        ln.tag = qn(tag)
        tcPr.append(ln)

    _add_border("a:lnT", top)
    _add_border("a:lnB", bottom)
    _add_border("a:lnL", left)
    _add_border("a:lnR", right)


def set_cell_fill(cell, hex_color: str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(el)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", hex_color.lstrip("#"))


def add_textbox(slide, x, y, cx, cy, text, font_name, font_size_pt,
                hex_color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(hex_color)
    return txBox


def add_rect(slide, x, y, cx, cy, hex_fill, line_color=None, line_width_emu=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(x), Emu(y), Emu(cx), Emu(cy)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_fill)
    line = shape.line
    if line_color:
        line.color.rgb = hex_to_rgb(line_color)
        line.width = Emu(line_width_emu)
    else:
        line.fill.background()
    return shape


def add_slide(prs) -> object:
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)


def set_slide_background(slide, hex_color: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def set_slide_note(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ---------------------------------------------------------------------------
# 공통 요소
# ---------------------------------------------------------------------------

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "outputs", "logo_src.png")
LOGO_PATH = os.path.normpath(LOGO_PATH)


def add_header_table(slide, section_label: str = "{{section}}"):
    """헤더 TABLE: 1행 2열, fill=#627365, 텍스트 white."""
    x, y, cx, cy = 614160, 441327, 10944225, 650267
    table = slide.shapes.add_table(1, 2, Emu(x), Emu(y), Emu(cx), Emu(cy)).table

    # 열 폭 배분: 좌 40% / 우 60%
    table.columns[0].width = Emu(int(cx * 0.40))
    table.columns[1].width = Emu(int(cx * 0.60))

    border_spec = (12700, "#000000")  # 1pt black

    for col_idx, text in enumerate([section_label, ""]):
        cell = table.cell(0, col_idx)
        set_cell_fill(cell, PAL["primary"])
        set_cell_border(
            cell,
            top=border_spec,
            bottom=border_spec,
            left=None,
            right=None,
        )
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = "Pretendard ExtraBold"
        run.font.size = Pt(19)
        run.font.color.rgb = hex_to_rgb(PAL["white"])


def add_head_message(slide, text: str = "{{head_message}}"):
    add_textbox(
        slide,
        x=628044, y=1141611, cx=10944225, cy=307777,
        text=text,
        font_name="Pretendard SemiBold",
        font_size_pt=14,
        hex_color=PAL["accent_warm"],
    )


def add_logo(slide):
    x, y, cx, cy = 623888, 6492899, 959322, 147158
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Emu(x), Emu(y), Emu(cx), Emu(cy))
    else:
        add_textbox(
            slide, x=x, y=y, cx=cx, cy=cy,
            text="SRC",
            font_name="Pretendard SemiBold",
            font_size_pt=9,
            hex_color=PAL["text_main"],
        )


def add_page_number(slide, number: int):
    add_textbox(
        slide,
        x=10424072, y=6482324, cx=1144041, cy=230832,
        text=str(number),
        font_name="Pretendard Light",
        font_size_pt=9,
        hex_color=PAL["page_num"],
        align=PP_ALIGN.RIGHT,
    )


def add_common_elements(slide, slide_number: int, section_label: str = "{{section}}"):
    add_header_table(slide, section_label)
    add_head_message(slide)
    add_logo(slide)
    add_page_number(slide, slide_number)


# ---------------------------------------------------------------------------
# 슬라이드 1: title_slide
# ---------------------------------------------------------------------------

def build_title_slide(prs):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["dark"])

    # 장식선 1
    add_rect(slide, x=0, y=1200000, cx=SLIDE_W, cy=8000, hex_fill=PAL["primary"])

    # 제목박스
    add_textbox(
        slide,
        x=1200000, y=2000000, cx=9800000, cy=1200000,
        text="{{title}}",
        font_name="Pretendard SemiBold",
        font_size_pt=32,
        hex_color=PAL["white"],
    )

    # 구분선
    add_rect(slide, x=1200000, y=3300000, cx=10000000, cy=6000, hex_fill=PAL["primary"])

    # 부제목
    add_textbox(
        slide,
        x=1200000, y=3400000, cx=9800000, cy=800000,
        text="{{subtitle}}",
        font_name="Pretendard Light",
        font_size_pt=13,
        hex_color=PAL["light_green"],
    )

    # 날짜
    add_textbox(
        slide,
        x=1200000, y=4600000, cx=9800000, cy=300000,
        text="{{date}}",
        font_name="Pretendard Light",
        font_size_pt=9,
        hex_color=PAL["light_green"],
    )

    add_logo(slide)
    set_slide_note(slide, "layout_name: title_slide")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 2: content_text
# ---------------------------------------------------------------------------

def build_content_text(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 본문
    body_text = "{{bullets}}"
    txBox = slide.shapes.add_textbox(Emu(628044), Emu(1454151), Emu(10973754), Emu(5050000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = body_text
    run.font.name = "Pretendard Light"
    run.font.size = Pt(9)
    run.font.color.rgb = hex_to_rgb(PAL["text_main"])

    # 출처
    add_textbox(
        slide,
        x=628044, y=6237289, cx=10944225, cy=200000,
        text="Source: {{source}}",
        font_name="Pretendard Light",
        font_size_pt=6,
        hex_color=PAL["page_num"],
    )

    set_slide_note(slide, "layout_name: content_text")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 3: content_chart
# ---------------------------------------------------------------------------

def build_content_chart(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 서브라벨 TABLE
    sub_table = slide.shapes.add_table(1, 1, Emu(614160), Emu(1449388), Emu(5292726), Emu(288000)).table
    cell = sub_table.cell(0, 0)
    set_cell_fill(cell, PAL["gold"])
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{chart_label}}"
    run.font.name = "Pretendard SemiBold"
    run.font.size = Pt(10)
    run.font.color.rgb = hex_to_rgb(PAL["white"])

    # 차트 플레이스홀더 박스 (name으로 식별)
    chart_box = add_rect(
        slide,
        x=614160, y=1773238, cx=10944225, cy=4462510,
        hex_fill=PAL["neutral_light"],
    )
    chart_box.name = "chart_placeholder"
    add_textbox(
        slide,
        x=614160, y=3900000, cx=10944225, cy=400000,
        text="[차트 영역]",
        font_name="Pretendard Light",
        font_size_pt=9,
        hex_color=PAL["page_num"],
        align=PP_ALIGN.CENTER,
    )

    # key_points 영역 (우측 오버레이)
    add_textbox(
        slide,
        x=8300000, y=1773238, cx=3900000, cy=4462510,
        text="{{key_points}}",
        font_name="Pretendard Light",
        font_size_pt=8,
        hex_color=PAL["text_main"],
    )

    set_slide_note(slide, "layout_name: content_chart")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 4: two_column_compare
# ---------------------------------------------------------------------------

def build_two_column_compare(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 좌열 헤더박스
    left_header = slide.shapes.add_table(1, 1, Emu(614160), Emu(1449388), Emu(5292726), Emu(288000)).table
    cell = left_header.cell(0, 0)
    set_cell_fill(cell, PAL["primary"])
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{left_title}}"
    run.font.name = "Pretendard ExtraBold"
    run.font.size = Pt(11)
    run.font.color.rgb = hex_to_rgb(PAL["white"])

    # 좌열 본문
    add_textbox(
        slide,
        x=628044, y=1754151, cx=5280000, cy=4050000,
        text="{{left_items}}",
        font_name="Pretendard Light",
        font_size_pt=9,
        hex_color=PAL["text_main"],
    )

    # 구분선
    add_rect(slide, x=6003000, y=1449388, cx=12700, cy=4700000, hex_fill=PAL["light_green"])

    # 우열 헤더박스
    right_header = slide.shapes.add_table(1, 1, Emu(6277940), Emu(1449388), Emu(5290172), Emu(288000)).table
    cell = right_header.cell(0, 0)
    set_cell_fill(cell, PAL["accent_blue"])
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{right_title}}"
    run.font.name = "Pretendard ExtraBold"
    run.font.size = Pt(11)
    run.font.color.rgb = hex_to_rgb(PAL["white"])

    # 우열 본문
    add_textbox(
        slide,
        x=6290000, y=1754151, cx=5280000, cy=4050000,
        text="{{right_items}}",
        font_name="Pretendard Light",
        font_size_pt=9,
        hex_color=PAL["text_main"],
    )

    set_slide_note(slide, "layout_name: two_column_compare")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 5: three_column_summary
# ---------------------------------------------------------------------------

def build_three_column_summary(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    col_w = 3648075
    x_starts = [614160, 4262235, 7910310]
    col_labels = ["{{col1_title}}", "{{col2_title}}", "{{col3_title}}"]
    col_bullets = [
        "{{col1_items}}",
        "{{col2_items}}",
        "{{col3_items}}",
    ]

    for i, (x0, label, bullets) in enumerate(zip(x_starts, col_labels, col_bullets)):
        # 상단 구분선
        add_rect(slide, x=x0, y=1454151, cx=col_w - 50000, cy=15000, hex_fill=PAL["primary"])

        # 열 제목
        add_textbox(
            slide,
            x=x0, y=1490000, cx=col_w - 50000, cy=250000,
            text=label,
            font_name="Pretendard SemiBold",
            font_size_pt=11,
            hex_color=PAL["primary"],
        )

        # 항목 bullets
        txBox = slide.shapes.add_textbox(Emu(x0), Emu(1760000), Emu(col_w - 50000), Emu(4694151))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullets  # just the marker
        run.font.name = "Pretendard Light"
        run.font.size = Pt(9)
        run.font.color.rgb = hex_to_rgb(PAL["text_main"])

    set_slide_note(slide, "layout_name: three_column_summary")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 6: table_slide
# ---------------------------------------------------------------------------

def build_table_slide(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    rows = 4
    cols = 3
    tbl = slide.shapes.add_table(
        rows, cols,
        Emu(614160), Emu(1773238),
        Emu(10944225), Emu(4462510)
    ).table

    headers = ["{{header_1}}", "{{header_2}}", "{{header_3}}"]
    data = [
        ["{{row1_col1}}", "{{row1_col2}}", "{{row1_col3}}"],
        ["{{row2_col1}}", "{{row2_col2}}", "{{row2_col3}}"],
        ["{{row3_col1}}", "{{row3_col2}}", "{{row3_col3}}"],
    ]

    border_spec = (12700, "#000000")

    for col_idx, header_text in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        set_cell_fill(cell, PAL["primary"])
        set_cell_border(cell, top=border_spec, bottom=border_spec, left=None, right=None)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header_text
        run.font.name = "Pretendard"
        run.font.size = Pt(9)
        run.font.color.rgb = hex_to_rgb(PAL["white"])

    for row_idx, row_data in enumerate(data):
        fill_color = PAL["neutral_light"] if row_idx % 2 == 0 else PAL["white"]
        for col_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            set_cell_fill(cell, fill_color)
            set_cell_border(cell, top=None, bottom=None, left=None, right=None)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Pretendard Light"
            run.font.size = Pt(9)
            run.font.color.rgb = hex_to_rgb(PAL["text_main"])

    set_slide_note(slide, "layout_name: table_slide")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 7: roadmap_timeline
# ---------------------------------------------------------------------------

def build_roadmap_timeline(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 타임라인 기준선
    timeline_y = 3500000
    add_rect(slide, x=800000, y=timeline_y, cx=10600000, cy=12700, hex_fill=PAL["primary"])

    node_count = 5
    node_size = 400000
    node_colors = [PAL["primary"], PAL["accent_warm"], PAL["primary"], PAL["accent_warm"], PAL["primary"]]
    x_positions = [
        800000 + int(i * 10600000 / (node_count - 1)) - node_size // 2
        for i in range(node_count)
    ]
    node_y = timeline_y - node_size // 2

    period_labels = ["{{period_1}}", "{{period_2}}", "{{period_3}}", "{{period_4}}", "{{period_5}}"]
    node_titles = ["{{node_title_1}}", "{{node_title_2}}", "{{node_title_3}}", "{{node_title_4}}", "{{node_title_5}}"]

    for i, (x0, color) in enumerate(zip(x_positions, node_colors)):
        # OVAL 노드
        oval = slide.shapes.add_shape(
            9,  # MSO_SHAPE_TYPE.OVAL
            Emu(x0), Emu(node_y), Emu(node_size), Emu(node_size)
        )
        oval.fill.solid()
        oval.fill.fore_color.rgb = hex_to_rgb(color)
        oval.line.fill.background()

        # 기간 레이블 (위 or 아래 지그재그)
        if i % 2 == 0:
            label_y = node_y - 400000
            title_y = node_y + node_size + 50000
        else:
            label_y = node_y + node_size + 50000
            title_y = node_y - 400000

        label_y = max(label_y, 0)
        title_y = max(title_y, 0)

        add_textbox(
            slide,
            x=x0 - 200000, y=label_y, cx=node_size + 400000, cy=300000,
            text=period_labels[i],
            font_name="Pretendard SemiBold",
            font_size_pt=9,
            hex_color=PAL["primary"],
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            x=x0 - 300000, y=title_y, cx=node_size + 600000, cy=300000,
            text=node_titles[i],
            font_name="Pretendard SemiBold",
            font_size_pt=8,
            hex_color=PAL["text_main"],
            align=PP_ALIGN.CENTER,
        )

    set_slide_note(slide, "layout_name: roadmap_timeline")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 8: closing_slide
# ---------------------------------------------------------------------------

def build_closing_slide(prs):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["dark"])

    # 장식선
    add_rect(slide, x=0, y=1800000, cx=SLIDE_W, cy=8000, hex_fill=PAL["primary"])

    # 제목
    add_textbox(
        slide,
        x=1200000, y=1900000, cx=9800000, cy=1000000,
        text="{{closing_title}}",
        font_name="Pretendard SemiBold",
        font_size_pt=22,
        hex_color=PAL["white"],
    )

    # 메시지
    add_textbox(
        slide,
        x=1200000, y=3000000, cx=9800000, cy=700000,
        text="{{closing_message}}",
        font_name="Pretendard Light",
        font_size_pt=11,
        hex_color=PAL["light_green"],
    )

    # takeaways
    add_textbox(
        slide,
        x=1200000, y=3800000, cx=9800000, cy=1800000,
        text="{{takeaways}}",
        font_name="Pretendard Light",
        font_size_pt=10,
        hex_color=PAL["light_green"],
    )

    add_logo(slide)
    set_slide_note(slide, "layout_name: closing_slide")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 9: section_divider
# ---------------------------------------------------------------------------

def build_section_divider(prs):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["dark"])

    # 좌측 강조 바
    add_rect(slide, x=800000, y=1800000, cx=50000, cy=3000000, hex_fill=PAL["primary"])

    # 섹션 번호 (대형)
    add_textbox(
        slide,
        x=1000000, y=1700000, cx=4000000, cy=1600000,
        text="{{section_number}}",
        font_name="Pretendard ExtraBold",
        font_size_pt=96,
        hex_color=PAL["primary"],
        bold=True,
    )

    # 구분선
    add_rect(slide, x=1000000, y=3400000, cx=8000000, cy=6000, hex_fill=PAL["gold"])

    # 섹션 제목
    add_textbox(
        slide,
        x=1000000, y=3480000, cx=10000000, cy=700000,
        text="{{section_title}}",
        font_name="Pretendard SemiBold",
        font_size_pt=24,
        hex_color=PAL["white"],
    )

    # 섹션 부제목
    add_textbox(
        slide,
        x=1000000, y=4300000, cx=10000000, cy=500000,
        text="{{section_subtitle}}",
        font_name="Pretendard Light",
        font_size_pt=11,
        hex_color=PAL["light_green"],
    )

    add_logo(slide)
    set_slide_note(slide, "layout_name: section_divider")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 10: kpi_metrics
# ---------------------------------------------------------------------------

def build_kpi_metrics(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 4개 KPI 카드 — 콘텐츠 영역(x=614160, w=10944225)에 균등 배치
    card_w = 2500000
    gap = (10944225 - 4 * card_w) // 3  # ≈ 314742
    card_y = 1950000
    card_h = 4100000

    for i in range(4):
        n = i + 1
        x0 = 614160 + i * (card_w + gap)

        # 카드 배경
        add_rect(slide, x=x0, y=card_y, cx=card_w, cy=card_h,
                 hex_fill=PAL["neutral_light"])

        # 상단 액센트 바
        add_rect(slide, x=x0, y=card_y, cx=card_w, cy=80000,
                 hex_fill=PAL["primary"])

        # 대형 수치
        add_textbox(
            slide,
            x=x0 + 80000, y=card_y + 200000, cx=card_w - 160000, cy=1600000,
            text=f"{{{{kpi_{n}_value}}}}",
            font_name="Pretendard SemiBold",
            font_size_pt=36,
            hex_color=PAL["primary"],
            align=PP_ALIGN.CENTER,
        )

        # 얇은 구분선
        add_rect(slide, x=x0 + 200000, y=card_y + 1900000,
                 cx=card_w - 400000, cy=6000, hex_fill=PAL["light_green"])

        # 레이블
        add_textbox(
            slide,
            x=x0 + 80000, y=card_y + 1980000, cx=card_w - 160000, cy=600000,
            text=f"{{{{kpi_{n}_label}}}}",
            font_name="Pretendard SemiBold",
            font_size_pt=10,
            hex_color=PAL["text_main"],
            align=PP_ALIGN.CENTER,
        )

        # 보조 노트
        add_textbox(
            slide,
            x=x0 + 80000, y=card_y + 2650000, cx=card_w - 160000, cy=500000,
            text=f"{{{{kpi_{n}_note}}}}",
            font_name="Pretendard Light",
            font_size_pt=8,
            hex_color=PAL["page_num"],
            align=PP_ALIGN.CENTER,
        )

    set_slide_note(slide, "layout_name: kpi_metrics")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 11: image_gallery
# ---------------------------------------------------------------------------

def build_image_gallery(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 2행 × 3열 이미지 그리드
    img_w = 3448075   # (10944225 / 3) ≈ 3648075 - gap
    img_h = 2100000
    label_h = 220000
    gap_x = 150000
    gap_y = 130000
    start_x = 614160
    start_y = 1540000

    for row in range(2):
        for col in range(3):
            n = row * 3 + col + 1
            x0 = start_x + col * (img_w + gap_x)
            y0 = start_y + row * (img_h + label_h + gap_y)

            # 이미지 플레이스홀더 박스
            img_box = add_rect(
                slide, x=x0, y=y0, cx=img_w, cy=img_h,
                hex_fill=PAL["neutral_light"],
            )
            img_box.name = f"image_placeholder_{n}"

            # 이미지 번호 인디케이터
            add_textbox(
                slide,
                x=x0, y=y0 + img_h // 2 - 120000, cx=img_w, cy=240000,
                text=f"[이미지 {n}]",
                font_name="Pretendard Light",
                font_size_pt=8,
                hex_color=PAL["page_num"],
                align=PP_ALIGN.CENTER,
            )

            # 라벨
            add_textbox(
                slide,
                x=x0, y=y0 + img_h + 10000, cx=img_w, cy=label_h,
                text=f"{{{{img_{n}_label}}}}",
                font_name="Pretendard SemiBold",
                font_size_pt=8,
                hex_color=PAL["text_main"],
                align=PP_ALIGN.CENTER,
            )

    set_slide_note(slide, "layout_name: image_gallery")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 12: table_chart_combo
# ---------------------------------------------------------------------------

def build_table_chart_combo(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 차트 서브라벨 (우측 상단)
    sub_table = slide.shapes.add_table(
        1, 1, Emu(6277940), Emu(1449388), Emu(5290172), Emu(288000)
    ).table
    cell = sub_table.cell(0, 0)
    set_cell_fill(cell, PAL["gold"])
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{chart_label}}"
    run.font.name = "Pretendard SemiBold"
    run.font.size = Pt(10)
    run.font.color.rgb = hex_to_rgb(PAL["white"])

    # 표 (좌측)
    rows = 4
    cols = 3
    tbl = slide.shapes.add_table(
        rows, cols,
        Emu(614160), Emu(1773238),
        Emu(5290172), Emu(4462510)
    ).table

    headers = ["{{header_1}}", "{{header_2}}", "{{header_3}}"]
    data = [
        ["{{row1_col1}}", "{{row1_col2}}", "{{row1_col3}}"],
        ["{{row2_col1}}", "{{row2_col2}}", "{{row2_col3}}"],
        ["{{row3_col1}}", "{{row3_col2}}", "{{row3_col3}}"],
    ]

    border_spec = (12700, "#000000")

    for col_idx, header_text in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        set_cell_fill(cell, PAL["primary"])
        set_cell_border(cell, top=border_spec, bottom=border_spec, left=None, right=None)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header_text
        run.font.name = "Pretendard"
        run.font.size = Pt(9)
        run.font.color.rgb = hex_to_rgb(PAL["white"])

    for row_idx, row_data in enumerate(data):
        fill_color = PAL["neutral_light"] if row_idx % 2 == 0 else PAL["white"]
        for col_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            set_cell_fill(cell, fill_color)
            set_cell_border(cell, top=None, bottom=None, left=None, right=None)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Pretendard Light"
            run.font.size = Pt(9)
            run.font.color.rgb = hex_to_rgb(PAL["text_main"])

    # 차트 플레이스홀더 박스 (우측, name으로 식별)
    chart_box = add_rect(
        slide,
        x=6277940, y=1773238, cx=5290172, cy=4462510,
        hex_fill=PAL["neutral_light"],
    )
    chart_box.name = "chart_placeholder"
    add_textbox(
        slide,
        x=6277940, y=3900000, cx=5290172, cy=400000,
        text="[차트 영역]",
        font_name="Pretendard Light",
        font_size_pt=9,
        hex_color=PAL["page_num"],
        align=PP_ALIGN.CENTER,
    )

    set_slide_note(slide, "layout_name: table_chart_combo")
    return slide


# ---------------------------------------------------------------------------
# 슬라이드 13: zone_base — 존 조립용 빈 콘텐츠 슬라이드
# 헤더/푸터 크롬만 가지고 콘텐츠 영역은 비어 있다.
# assembler.py 의 _post_process_zones() 가 이 슬라이드에 존 컴포넌트를 삽입한다.
# ---------------------------------------------------------------------------

def build_zone_base(prs, slide_number: int):
    slide = add_slide(prs)
    set_slide_background(slide, PAL["white"])
    add_common_elements(slide, slide_number)

    # 출처 플레이스홀더 (비어 있으면 후처리에서 제거됨)
    add_textbox(
        slide,
        x=628044, y=6237289, cx=10944225, cy=200000,
        text="Source: {{source}}",
        font_name="Pretendard Light",
        font_size_pt=6,
        hex_color=PAL["page_num"],
    )

    set_slide_note(slide, "layout_name: zone_base")
    return slide


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------

def main():
    base_dir = os.path.normpath(
        os.path.join(os.path.dirname(__file__), "..", "..")
    )
    output_pptx = os.path.join(base_dir, "outputs", "component_template_report.pptx")
    output_index = os.path.join(base_dir, "outputs", "component_template_report_index.json")

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    build_title_slide(prs)              # 슬라이드 1
    build_content_text(prs, 2)          # 슬라이드 2
    build_content_chart(prs, 3)         # 슬라이드 3
    build_two_column_compare(prs, 4)    # 슬라이드 4
    build_three_column_summary(prs, 5)  # 슬라이드 5
    build_table_slide(prs, 6)           # 슬라이드 6
    build_roadmap_timeline(prs, 7)      # 슬라이드 7
    build_closing_slide(prs)            # 슬라이드 8
    build_section_divider(prs)          # 슬라이드 9
    build_kpi_metrics(prs, 10)          # 슬라이드 10
    build_image_gallery(prs, 11)        # 슬라이드 11
    build_table_chart_combo(prs, 12)    # 슬라이드 12
    build_zone_base(prs, 13)            # 슬라이드 13 (존 조립 베이스)

    prs.save(output_pptx)
    print(f"PPTX 저장 완료: {output_pptx}")

    index = {
        "template": output_pptx,
        "layouts": {
            "title_slide": 1,
            "content_text": 2,
            "content_chart": 3,
            "two_column_compare": 4,
            "three_column_summary": 5,
            "table_slide": 6,
            "roadmap_timeline": 7,
            "closing_slide": 8,
            "section_divider": 9,
            "kpi_metrics": 10,
            "image_gallery": 11,
            "table_chart_combo": 12,
            "zone_base": 13,
            # zone 으로 선언된 슬라이드도 zone_base 템플릿을 사용
            "zone": 13,
        },
    }
    with open(output_index, "w", encoding="utf-8") as f:
        json.dump(index, f, ensure_ascii=False, indent=2)
    print(f"인덱스 저장 완료: {output_index}")


if __name__ == "__main__":
    main()
