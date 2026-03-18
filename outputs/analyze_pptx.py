import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

def analyze_pptx(filepath, label):
    print('=' * 80)
    print(f'FILE: {label}')
    print('=' * 80)

    prs = Presentation(filepath)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f'Slide size: {slide_width} x {slide_height} EMU  ({slide_width/914400:.3f} x {slide_height/914400:.3f} inches)')
    print()

    SHAPE_TYPES = {
        1:'AUTO_SHAPE', 3:'CHART', 5:'FREEFORM', 6:'GROUP', 9:'LINE',
        11:'LINKED_PICTURE', 13:'PICTURE', 14:'PLACEHOLDER', 16:'TABLE',
        17:'TEXT_BOX', 18:'MEDIA', 19:'TABLE', 24:'SMART_ART'
    }

    def get_texts(shape):
        texts = []
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = ''.join(run.text for run in para.runs)
                if line.strip():
                    texts.append(line)
        if shape.shape_type == 19:  # TABLE
            for row in shape.table.rows:
                for cell in row.cells:
                    ct = cell.text_frame.text.strip()
                    if ct:
                        texts.append(f'[CELL] {ct}')
        return texts

    def is_offcanvas(shape, sw, sh):
        left = shape.left if shape.left is not None else 0
        top = shape.top if shape.top is not None else 0
        return left < 0 or top < 0 or left > sw or top > sh

    # ── TASK 1 ──
    print('=' * 80)
    print('TASK 1: ALL TEXT CONTENT BY SLIDE')
    print('=' * 80)

    for slide_num, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else 'N/A'
        print(f'\n--- SLIDE {slide_num} (layout: {layout_name}) ---')
        for shape in slide.shapes:
            stype = SHAPE_TYPES.get(shape.shape_type, str(shape.shape_type))
            offcanvas = is_offcanvas(shape, slide_width, slide_height)
            flag = ' [OFF-CANVAS]' if offcanvas else ''
            texts = get_texts(shape)
            print(f'  Shape: "{shape.name}" | type={stype}{flag}')
            print(f'    pos: L={shape.left} T={shape.top} W={shape.width} H={shape.height}')
            if texts:
                for t in texts:
                    print(f'    TEXT: {t}')
            else:
                print(f'    (no text)')

    # ── TASK 2 ──
    print()
    print('=' * 80)
    print('TASK 2: OFF-CANVAS SHAPES SUMMARY')
    print('=' * 80)

    found = False
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if is_offcanvas(shape, slide_width, slide_height):
                found = True
                stype = SHAPE_TYPES.get(shape.shape_type, str(shape.shape_type))
                texts = get_texts(shape)
                print(f'\n  SLIDE {slide_num} | "{shape.name}" | type={stype}')
                print(f'    L={shape.left} T={shape.top} W={shape.width} H={shape.height}')
                if texts:
                    for t in texts:
                        print(f'    TEXT: {t}')
                else:
                    print(f'    (no text)')
    if not found:
        print('  None found in slides.')

    print()
    print('--- Checking MASTER and LAYOUTS ---')
    master_found = False
    for mi, master in enumerate(prs.slide_masters):
        for shape in master.shapes:
            if is_offcanvas(shape, slide_width, slide_height):
                master_found = True
                stype = SHAPE_TYPES.get(shape.shape_type, str(shape.shape_type))
                texts = get_texts(shape)
                print(f'  MASTER[{mi}] | "{shape.name}" | type={stype}')
                print(f'    L={shape.left} T={shape.top}')
                for t in texts:
                    print(f'    TEXT: {t}')
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                if is_offcanvas(shape, slide_width, slide_height):
                    master_found = True
                    stype = SHAPE_TYPES.get(shape.shape_type, str(shape.shape_type))
                    texts = get_texts(shape)
                    print(f'  LAYOUT "{layout.name}" | "{shape.name}" | type={stype}')
                    print(f'    L={shape.left} T={shape.top}')
                    for t in texts:
                        print(f'    TEXT: {t}')
    if not master_found:
        print('  None found in master/layouts.')


analyze_pptx('C:/Side project/ppt generator/outputs/template_test1.pptx', 'template_test1.pptx')
print()
print()
analyze_pptx('C:/Side project/ppt generator/outputs/template_test1 - \ub3b5\uc0ac\ubcf8.pptx', 'template_test1 - 복사본.pptx')
